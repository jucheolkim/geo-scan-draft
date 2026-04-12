import streamlit as st
import anthropic
import io
import json
import openpyxl
from datetime import datetime

from docx import Document
from docx.shared import RGBColor, Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt as PPTXPt, Cm as PPTXCm
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.text import PP_ALIGN


# ─────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────
def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def brand_rgb():
    return hex_to_rgb(st.session_state.brand_color)

def check_mention(text, brand_name):
    if not text.strip():
        return None
    keywords = [brand_name, brand_name[:2], brand_name.replace(' ', '')]
    return any(kw in text for kw in keywords)


# ─────────────────────────────────────────────
# Excel 읽기
# ─────────────────────────────────────────────
def load_from_excel(uploaded_file):
    wb = openpyxl.load_workbook(uploaded_file)
    sheets = wb.sheetnames
    data = {
        'brand_name': '', 'brand_color': '#4A90D9',
        'brand_category': '', 'brand_usp': '',
        'brand_target': '', 'brand_competitors': '',
        'brand_negative': '', 'brand_focus': '',
        'questions': [], 'answers': {
            'off': {i: '' for i in range(1, 8)},
            'on':  {i: '' for i in range(1, 8)},
            'gem': {i: '' for i in range(1, 8)},
        },
        'analysis_result': '',
        'cw_insights': [''] * 7,
        'overall_diagnosis': '',
        'priority_actions': '',
    }

    # ── 브랜드 정보 ──
    if '브랜드 정보' in sheets:
        ws = wb['브랜드 정보']
        info_map = {
            '브랜드명':    'brand_name',
            '카테고리':    'brand_category',
            '핵심 USP':   'brand_usp',
            '주요 타겟':   'brand_target',
            '경쟁 브랜드': 'brand_competitors',
            '부정 이미지': 'brand_negative',
            '강조 포인트': 'brand_focus',
        }
        for row in ws.iter_rows(min_row=2, values_only=True):
            if not row[0]:
                continue
            label, value = row[0], row[1]
            if label in info_map and value:
                data[info_map[label]] = str(value)

    # ── 질문 + 답변 ──
    if 'AI 답변 수집' in sheets:
        ws2 = wb['AI 답변 수집']
        questions = []
        for row in ws2.iter_rows(min_row=2, values_only=True):
            if not row[0] or not str(row[0]).startswith('Q'):
                continue
            n = int(str(row[0]).replace('Q', ''))
            q_dict = {
                'question':    str(row[1]) if row[1] else '',
                'type':        str(row[2]) if row[2] else '',
                'stage':       str(row[3]) if row[3] else '',
                'check_point': str(row[4]) if len(row) > 4 and row[4] else '',
                'data':        [],
            }
            questions.append(q_dict)
            if len(row) >= 8 and row[5] is not None:
                data['answers']['off'][n] = str(row[5]) if row[5] else ''
                data['answers']['on'][n]  = str(row[6]) if row[6] else ''
                data['answers']['gem'][n] = str(row[7]) if row[7] else ''
            else:
                data['answers']['off'][n] = str(row[4]) if len(row) > 4 and row[4] else ''
                data['answers']['on'][n]  = str(row[5]) if len(row) > 5 and row[5] else ''
                data['answers']['gem'][n] = str(row[6]) if len(row) > 6 and row[6] else ''
        data['questions'] = questions

    # ── Claude 분석 결과 ──
    if 'Claude 분석' in sheets:
        ws4 = wb['Claude 분석']
        for row in ws4.iter_rows(min_row=2, values_only=True):
            if row[0]:
                data['analysis_result'] = str(row[0])
                data['overall_diagnosis'] = str(row[0])[:300]
                break

    return data


# ─────────────────────────────────────────────
# [1] 질문지 Word 생성
# ─────────────────────────────────────────────
def create_question_word(d):
    from docx.oxml.ns import qn as docx_qn

    FONT = "맑은 고딕"
    br, bg, bb = hex_to_rgb(d['brand_color'])
    cr, cg, cb = (112, 48, 160)
    BRAND_HEX = d['brand_color'].replace('#', '')

    lr = int(br + (255 - br) * 0.5)
    lg = int(bg + (255 - bg) * 0.5)
    lb = int(bb + (255 - bb) * 0.5)
    BRAND_LIGHT_HEX = f'{lr:02X}{lg:02X}{lb:02X}'
    CW_HEX    = 'EADCF4'
    GRAY_HEX  = 'F7F7F7'
    WHITE_HEX = 'FFFFFF'
    GRAY2_HEX = 'FAFAFA'

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Cm(1.8)
        section.bottom_margin = Cm(1.8)
        section.left_margin   = Cm(2.2)
        section.right_margin  = Cm(2.2)

    def set_bg(cell, hex_color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(docx_qn('w:val'), 'clear')
        shd.set(docx_qn('w:color'), 'auto')
        shd.set(docx_qn('w:fill'), hex_color)
        tcPr.append(shd)

    def set_cell_margins(cell, top=60, bottom=60, left=120, right=120):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcMar = OxmlElement('w:tcMar')
        for side, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
            node = OxmlElement(f'w:{side}')
            node.set(docx_qn('w:w'), str(val))
            node.set(docx_qn('w:type'), 'dxa')
            tcMar.append(node)
        tcPr.append(tcMar)

    def r(para, text, size=None, bold=False, color=None, font=FONT):
        run = para.add_run(text)
        run.font.name = font
        if size:
            run.font.size = Pt(size)
        run.bold = bold
        if color:
            try:
                run.font.color.rgb = RGBColor(int(color[0]), int(color[1]), int(color[2]))
            except Exception:
                pass
        return run

    def add_border_line(doc, color_hex, thickness=4):
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bot = OxmlElement('w:bottom')
        bot.set(docx_qn('w:val'), 'single')
        bot.set(docx_qn('w:sz'), str(thickness))
        bot.set(docx_qn('w:space'), '1')
        bot.set(docx_qn('w:color'), color_hex)
        pBdr.append(bot)
        pPr.append(pBdr)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(2)
        return p

    # 표지
    p0 = doc.add_paragraph()
    p0.paragraph_format.space_after = Pt(2)
    r(p0, "CREAMWORKS  ×", size=14, color=(124, 92, 191))
    add_border_line(doc, BRAND_HEX, thickness=4)
    p1 = doc.add_paragraph()
    p1.paragraph_format.space_before = Pt(6)
    p1.paragraph_format.space_after  = Pt(4)
    r(p1, d['brand_name'], size=32, color=(26, 23, 20))
    p2 = doc.add_paragraph()
    p2.paragraph_format.space_after = Pt(4)
    r(p2, "AI 진단 질문지", size=18, color=(cr, cg, cb))
    p3 = doc.add_paragraph()
    p3.paragraph_format.space_after = Pt(30)
    r(p3, f"Presented by CREAMWORKS  ·  {datetime.now().strftime('%Y.%m')}", color=(117, 117, 117))
    doc.add_page_break()

    # 실행 전 필수 세팅
    p_h = doc.add_paragraph()
    p_h.paragraph_format.space_after = Pt(8)
    r(p_h, "실행 전 필수 세팅", size=12)

    setup_t = doc.add_table(rows=2, cols=2)
    setup_t.style = 'Table Grid'
    setup_data = [
        ("ChatGPT",
         "① 메모리 OFF → 검색 OFF  (새 채팅 시작 후 메모리·검색 모두 비활성화)\n"
         "② 메모리 OFF → 검색 ON  (새 채팅 시작 후 검색만 활성화)"),
        ("Gemini",
         "③ 시크릿 모드 → gemini.google.com 접속  (로그아웃 상태에서 질문)"),
    ]
    for ri, (label, content) in enumerate(setup_data):
        set_bg(setup_t.rows[ri].cells[0], CW_HEX)
        set_bg(setup_t.rows[ri].cells[1], GRAY_HEX)
        set_cell_margins(setup_t.rows[ri].cells[0])
        set_cell_margins(setup_t.rows[ri].cells[1])
        r(setup_t.rows[ri].cells[0].paragraphs[0], label, color=(26, 23, 20))
        r(setup_t.rows[ri].cells[1].paragraphs[0], content, size=9.5, color=(26, 23, 20))

    doc.add_paragraph().paragraph_format.space_after = Pt(4)
    notice_t = doc.add_table(rows=1, cols=1)
    notice_t.style = 'Table Grid'
    set_bg(notice_t.rows[0].cells[0], CW_HEX)
    set_cell_margins(notice_t.rows[0].cells[0])
    np_ = notice_t.rows[0].cells[0].paragraphs[0]
    r(np_, "📌  각 질문은 새 채팅에서 입력하지 않고, 동일한 채팅 내에서 Q1→Q7 순서대로 연속 입력합니다.\n", size=9.5, color=(85, 85, 85))
    r(np_, "📌  답변은 복사해서 별도 파일에 Q번호와 함께 저장해주세요. (예: Q1_GPT검색OFF.txt)", size=9.5, color=(85, 85, 85))

    doc.add_paragraph().paragraph_format.space_after = Pt(8)
    p_h2 = doc.add_paragraph()
    p_h2.paragraph_format.space_after = Pt(8)
    r(p_h2, "진단 질문 7개", size=16)

    # Q1~Q7
    for i, q_data in enumerate(d['questions']):
        n     = i + 1
        q_txt = q_data.get('question', '')
        qtype = q_data.get('type', '')
        check = q_data.get('check_point', '')
        dlist = q_data.get('data', [])
        clean_check = check.replace('확인 포인트:', '').replace('확인 포인트 :', '').strip()

        pq = doc.add_paragraph()
        pq.paragraph_format.space_before = Pt(8)
        pq.paragraph_format.space_after  = Pt(3)
        r(pq, f"Q{n}.  ", size=14, color=(br, bg, bb))
        r(pq, q_txt, size=14, color=(26, 23, 20))

        type_t = doc.add_table(rows=1, cols=2)
        type_t.style = 'Table Grid'
        set_bg(type_t.rows[0].cells[0], BRAND_LIGHT_HEX)
        set_bg(type_t.rows[0].cells[1], GRAY_HEX)
        set_cell_margins(type_t.rows[0].cells[0])
        set_cell_margins(type_t.rows[0].cells[1])
        tp = type_t.rows[0].cells[0].paragraphs[0]
        r(tp, "유형", color=(85, 85, 85))
        tp.add_run('\n')
        r(tp, qtype, color=(26, 23, 20))
        cp_ = type_t.rows[0].cells[1].paragraphs[0]
        r(cp_, "확인 포인트 : ", color=(85, 85, 85))
        r(cp_, clean_check, color=(26, 23, 20))

        pd_ = doc.add_paragraph()
        pd_.paragraph_format.space_before = Pt(6)
        pd_.paragraph_format.space_after  = Pt(2)
        r(pd_, "📊  선정 근거 데이터", color=(85, 85, 85))

        data_rows = max(len(dlist), 1)
        dt = doc.add_table(rows=data_rows + 1, cols=3)
        dt.style = 'Table Grid'
        for ci, h_txt in enumerate(["출처", "데이터", "연도"]):
            set_bg(dt.rows[0].cells[ci], CW_HEX)
            set_cell_margins(dt.rows[0].cells[ci])
            r(dt.rows[0].cells[ci].paragraphs[0], h_txt, color=(85, 85, 85))
        for j, ditem in enumerate(dlist):
            row_bg = WHITE_HEX if j % 2 == 0 else GRAY2_HEX
            for ci in range(3):
                set_bg(dt.rows[j+1].cells[ci], row_bg)
                set_cell_margins(dt.rows[j+1].cells[ci])
            r(dt.rows[j+1].cells[0].paragraphs[0], ditem.get('source', ''), color=(26, 23, 20))
            r(dt.rows[j+1].cells[1].paragraphs[0], ditem.get('content', ''), color=(26, 23, 20))
            r(dt.rows[j+1].cells[2].paragraphs[0], ditem.get('year', ''), color=(26, 23, 20))

        ins_t = doc.add_table(rows=1, cols=1)
        ins_t.style = 'Table Grid'
        set_bg(ins_t.rows[0].cells[0], CW_HEX)
        set_cell_margins(ins_t.rows[0].cells[0])
        ip = ins_t.rows[0].cells[0].paragraphs[0]
        r(ip, "→  ", color=(0, 0, 0))
        r(ip, clean_check if clean_check else f"이 질문에서 {d['brand_name']}이(가) 어떻게 언급되는지 확인하세요.", color=(0, 0, 0))
        doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # 질문 선정 근거 요약
    doc.add_page_break()
    p_sum = doc.add_paragraph()
    p_sum.paragraph_format.space_after = Pt(4)
    r(p_sum, "질문 선정 근거 요약", size=14)
    p_sum_desc = doc.add_paragraph()
    p_sum_desc.paragraph_format.space_after = Pt(8)
    r(p_sum_desc, "이 7개 질문은 다음 자료를 교차 분석해 도출했습니다.", size=9.5, color=(85, 85, 85))

    sum_t = doc.add_table(rows=1, cols=3)
    sum_t.style = 'Table Grid'
    for ci, h_txt in enumerate(["자료명", "발행처", "연도"]):
        set_bg(sum_t.rows[0].cells[ci], CW_HEX)
        set_cell_margins(sum_t.rows[0].cells[ci])
        r(sum_t.rows[0].cells[ci].paragraphs[0], h_txt, size=9, color=(85, 85, 85))

    all_sources = {}
    for q_data in d['questions']:
        for ditem in q_data.get('data', []):
            src = ditem.get('source', '').strip()
            if src and src not in all_sources:
                all_sources[src] = (ditem.get('content', ''), ditem.get('year', ''))
    for idx, (src, (content_val, yr)) in enumerate(all_sources.items()):
        row_bg = WHITE_HEX if idx % 2 == 0 else GRAY2_HEX
        new_row = sum_t.add_row()
        for ci in range(3):
            set_bg(new_row.cells[ci], row_bg)
            set_cell_margins(new_row.cells[ci])
        data_name = content_val[:40] + "…" if len(content_val) > 40 else content_val
        r(new_row.cells[0].paragraphs[0], data_name or src, size=9, color=(26, 23, 20))
        r(new_row.cells[1].paragraphs[0], src, size=9, color=(26, 23, 20))
        r(new_row.cells[2].paragraphs[0], yr, size=9, color=(26, 23, 20))

    doc.add_paragraph().paragraph_format.space_after = Pt(6)
    fp = doc.add_paragraph()
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r(fp, "CREAMWORKS  —  AI가 좋아하는 브랜드를 만듭니다", size=9, color=(85, 85, 85))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────
# [2] 전략제안서 Word 생성
# ─────────────────────────────────────────────
def create_proposal_word(d):
    from docx.oxml.ns import qn as docx_qn

    FONT = "맑은 고딕"
    br, bg, bb = hex_to_rgb(d['brand_color'])
    cr, cg, cb = (112, 48, 160)
    BRAND_HEX  = d['brand_color'].replace('#', '')
    CW_HEX     = 'EADCF4'
    GRAY_HEX   = 'F7F7F7'
    WHITE_HEX  = 'FFFFFF'
    GREEN_HEX  = 'D9F2D0'
    ORANGE_HEX = 'FAE2D5'
    DARK_HEX   = '1A1A1A'

    ai_keys   = ['off', 'on', 'gem']
    ai_labels = ['GPT 검색OFF', 'GPT 검색ON', 'Gemini']

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Cm(1.8)
        section.bottom_margin = Cm(1.8)
        section.left_margin   = Cm(2.2)
        section.right_margin  = Cm(2.2)

    def set_bg(cell, hex_color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_color)
        tcPr.append(shd)

    def r(para, text, size=10, bold=False, color=None, font=FONT):
        run = para.add_run(text)
        run.font.name = font
        run.font.size = Pt(size)
        run.bold = bold
        if color:
            try:
                run.font.color.rgb = RGBColor(int(color[0]), int(color[1]), int(color[2]))
            except Exception:
                pass
        return run

    def add_border_para(doc, color_hex, thickness=4):
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bot = OxmlElement('w:bottom')
        bot.set(qn('w:val'), 'single')
        bot.set(qn('w:sz'), str(thickness))
        bot.set(qn('w:space'), '1')
        bot.set(qn('w:color'), color_hex)
        pBdr.append(bot)
        pPr.append(pBdr)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(4)
        return p

    def part_header(doc, part_num, title):
        p = doc.add_paragraph()
        r(p, f"PART {part_num}  ", size=16, bold=True, color=(br, bg, bb))
        r(p, title, size=14, bold=True, color=(26, 23, 20))
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after  = Pt(8)
        return p

    # 표지
    p_logo = doc.add_paragraph()
    r(p_logo, "CREAMWORKS  |  GEO 컨설팅 제안서  |  Confidential", size=9, color=(85, 85, 85))
    p_logo.paragraph_format.space_after = Pt(2)
    add_border_para(doc, BRAND_HEX, 6)
    p_brand = doc.add_paragraph()
    p_brand.paragraph_format.space_before = Pt(40)
    p_brand.paragraph_format.space_after  = Pt(6)
    r(p_brand, d['brand_name'], size=36, color=(26, 23, 20))
    p_sub = doc.add_paragraph()
    r(p_sub, "AI 검색 최적화 (GEO) 전략 제안서", size=18, color=(cr, cg, cb))
    p_sub.paragraph_format.space_after = Pt(6)
    p_desc = doc.add_paragraph()
    r(p_desc, "실제 AI 답변 기반 현황 진단 + GEO 개선 전략", size=11, color=(85, 85, 85))
    p_desc.paragraph_format.space_after = Pt(40)
    doc.add_page_break()

    # 목적 박스
    obj_t = doc.add_table(rows=1, cols=1)
    obj_t.style = 'Table Grid'
    set_bg(obj_t.rows[0].cells[0], CW_HEX)
    op = obj_t.rows[0].cells[0].paragraphs[0]
    r(op, "📋  이 문서의 목적\n", size=10, bold=True, color=(cr, cg, cb))
    r(op, f"ChatGPT와 Gemini에 실제 소비자 질문 7개를 입력해 얻은 답변을 분석한 결과를 바탕으로, "
          f"{d['brand_name']}의 현재 AI 검색 노출 현황을 진단하고 구체적인 GEO 개선 전략을 제시합니다.",
      size=10, color=(26, 23, 20))
    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # PART 0: 질문 설계
    doc.add_page_break()
    part_header(doc, 0, "진단 질문 설계 — 근거와 방법론")
    q_table = doc.add_table(rows=len(d['questions']) + 1, cols=3)
    q_table.style = 'Table Grid'
    for ci, h in enumerate(["번호", "단계", "질문"]):
        set_bg(q_table.rows[0].cells[ci], CW_HEX)
        r(q_table.rows[0].cells[ci].paragraphs[0], h, size=9, bold=True, color=(85, 85, 85))
    for i, q_data in enumerate(d['questions']):
        n   = i + 1
        row = q_table.rows[n]
        bg  = WHITE_HEX if i % 2 == 0 else 'FAFAFA'
        for ci in range(3):
            set_bg(row.cells[ci], bg)
        r(row.cells[0].paragraphs[0], f"Q{n}", size=9)
        r(row.cells[1].paragraphs[0], q_data.get('stage', ''), size=9)
        r(row.cells[2].paragraphs[0], q_data.get('question', ''), size=9)
    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # PART 1: AI 진단 결과
    doc.add_page_break()
    part_header(doc, 1, "AI 진단 결과 — 현황 분석")
    conc_t = doc.add_table(rows=1, cols=1)
    conc_t.style = 'Table Grid'
    set_bg(conc_t.rows[0].cells[0], CW_HEX)
    cp = conc_t.rows[0].cells[0].paragraphs[0]
    r(cp, "🔑  핵심 결론\n", size=10, bold=True, color=(cr, cg, cb))
    analysis_preview = d['analysis_result'][:400] + "..." if len(d['analysis_result']) > 400 else d['analysis_result']
    r(cp, analysis_preview, size=10, color=(26, 23, 20))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # B2A 매트릭스
    mt = doc.add_table(rows=len(d['questions']) + 1, cols=5)
    mt.style = 'Table Grid'
    for ci, h in enumerate(["질문", "내용", "GPT 검색OFF", "GPT 검색ON", "Gemini"]):
        set_bg(mt.rows[0].cells[ci], CW_HEX)
        r(mt.rows[0].cells[ci].paragraphs[0], h, size=9, bold=True, color=(85, 85, 85))
    for i, q_data in enumerate(d['questions']):
        n   = i + 1
        row = mt.rows[n]
        bg  = WHITE_HEX if i % 2 == 0 else 'FEFFF0'
        for ci in range(5):
            set_bg(row.cells[ci], bg)
        r(row.cells[0].paragraphs[0], f"Q{n}", size=9)
        r(row.cells[1].paragraphs[0], q_data.get('question', '')[:30], size=9)
        for j, key in enumerate(ai_keys):
            ans     = d['answers'][key][n]
            mention = check_mention(ans, d['brand_name'])
            cell    = row.cells[j + 2]
            if ans.strip():
                if mention:
                    set_bg(cell, 'D9F2D0')
                    r(cell.paragraphs[0], "✅ 언급됨", size=9, color=(21, 87, 36))
                else:
                    set_bg(cell, 'FAE2D5')
                    r(cell.paragraphs[0], "❌ 미언급", size=9, color=(180, 50, 30))
            else:
                r(cell.paragraphs[0], "—", size=9, color=(150, 150, 150))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # PART 2: GEO 기초 진단
    doc.add_page_break()
    part_header(doc, 2, "브랜드 GEO 기초 진단")
    gd_t = doc.add_table(rows=1, cols=2)
    gd_t.style = 'Table Grid'
    set_bg(gd_t.rows[0].cells[0], GREEN_HEX)
    set_bg(gd_t.rows[0].cells[1], ORANGE_HEX)
    r(gd_t.rows[0].cells[0].paragraphs[0],
      f"✅ GEO 강점\n{d['brand_usp']}", size=9, color=(26, 23, 20))
    r(gd_t.rows[0].cells[1].paragraphs[0],
      f"❌ GEO 약점\n{d['brand_negative']}", size=9, color=(26, 23, 20))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    competitors = [c.strip() for c in d['brand_competitors'].split(',') if c.strip()]
    if competitors:
        comp_t = doc.add_table(rows=len(competitors) + 1, cols=3)
        comp_t.style = 'Table Grid'
        for ci, h in enumerate(["경쟁사", "AI 내 현재 포지션", f"{d['brand_name']}과의 격차"]):
            set_bg(comp_t.rows[0].cells[ci], CW_HEX)
            r(comp_t.rows[0].cells[ci].paragraphs[0], h, size=9, bold=True, color=(85, 85, 85))
        for i, comp in enumerate(competitors):
            row = comp_t.rows[i + 1]
            bg  = WHITE_HEX if i % 2 == 0 else 'FAFAFA'
            for ci in range(3):
                set_bg(row.cells[ci], bg)
            r(row.cells[0].paragraphs[0], comp, size=9)
            r(row.cells[1].paragraphs[0], "AI 내 포지션 분석 필요", size=9, color=(150, 150, 150))
            r(row.cells[2].paragraphs[0], "분석 결과 참조", size=9, color=(150, 150, 150))
        doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # PART 3: 전략 인사이트
    doc.add_page_break()
    part_header(doc, 3, "크림웍스 GEO 전략 인사이트")
    diag_t = doc.add_table(rows=1, cols=1)
    diag_t.style = 'Table Grid'
    set_bg(diag_t.rows[0].cells[0], CW_HEX)
    dp = diag_t.rows[0].cells[0].paragraphs[0]
    r(dp, "전체 현황 진단\n", size=10, bold=True, color=(cr, cg, cb))
    r(dp, d['overall_diagnosis'] or d['analysis_result'][:300], size=10, color=(26, 23, 20))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    for i, q_data in enumerate(d['questions']):
        n       = i + 1
        insight = d['cw_insights'][i] if i < len(d['cw_insights']) else ''
        if insight.strip():
            ins_t = doc.add_table(rows=2, cols=1)
            ins_t.style = 'Table Grid'
            set_bg(ins_t.rows[0].cells[0], DARK_HEX)
            r(ins_t.rows[0].cells[0].paragraphs[0],
              f"Q{n}. {q_data.get('question', '')}", size=10, bold=True, color=(br, bg, bb))
            set_bg(ins_t.rows[1].cells[0], CW_HEX)
            r(ins_t.rows[1].cells[0].paragraphs[0],
              f"💜 크림웍스 전략: {insight}", size=10, color=(cr, cg, cb))
            doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # PART 4: 액션 플랜
    doc.add_page_break()
    part_header(doc, 4, "우선 실행 액션 플랜")
    if d['priority_actions'].strip():
        for line in d['priority_actions'].split('\n'):
            if line.strip():
                p = doc.add_paragraph()
                r(p, line.strip(), size=10, color=(26, 23, 20))

    fp = doc.add_paragraph()
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r(fp, "CREAMWORKS  —  AI가 좋아하는 브랜드를 만듭니다", size=9, color=(150, 150, 150))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────
# [3] 질문지 PPT 생성
# ─────────────────────────────────────────────
def create_question_ppt(d):
    FONT_M = "맑은 고딕"
    FONT_R = "맑은 고딕"
    FONT_L = "맑은 고딕"
    br, bg, bb = hex_to_rgb(d['brand_color'])
    BRAND_RGB = PPTXRGBColor(br, bg, bb)
    CW_PURPLE = PPTXRGBColor(83, 39, 168)
    CW_LIGHT  = PPTXRGBColor(124, 92, 191)
    DARK      = PPTXRGBColor(26, 23, 20)
    GRAY      = PPTXRGBColor(85, 85, 85)
    WHITE     = PPTXRGBColor(255, 255, 255)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    blank = prs.slide_layouts[6]

    def add_textbox(slide, left, top, width, height, text,
                    font_name=FONT_L, font_size=16, bold=False,
                    color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(
            PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = PPTXPt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return tb

    def add_rect(slide, left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(
            1, PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    # 슬라이드 1: 표지
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 33.87, 19.05, PPTXRGBColor(245, 245, 245))
    add_rect(s1, 8.56, 4.55, 0.46, 5.39, CW_PURPLE)
    add_textbox(s1, 9.38, 4.73, 12.0, 1.81, "크림웍스\nGEO 컨설팅", FONT_M, 36, False, DARK)
    add_textbox(s1, 9.38, 6.55, 7.0, 1.27, "[AI 진단 설문지]", FONT_R, 24, False, DARK)
    add_textbox(s1, 9.38, 7.73, 14.1, 1.63,
                f"AI 검색시대, {d['brand_name']}이(가) ChatGPT·Gemini에서\n발견되고 추천되기 위한 진단 질문지입니다.",
                FONT_L, 16, False, DARK)
    add_textbox(s1, 11.20, 17.69, 11.47, 0.86,
                "CREAMWORKS  -  AI가 좋아하는 브랜드를 만듭니다",
                FONT_M, 14, False, CW_PURPLE)

    # 슬라이드 2: 실행 세팅
    s2 = prs.slides.add_slide(blank)
    add_textbox(s2, 1.25, 2.24, 7.32, 1.45, "실행 전 필수 세팅", FONT_R, 28, False, DARK)
    line = s2.shapes.add_shape(1, PPTXCm(1.25), PPTXCm(3.71), PPTXCm(31.37), PPTXCm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_RGB
    line.line.fill.background()
    add_rect(s2, 1.98, 5.46, 13.23, 7.19, PPTXRGBColor(255, 248, 240))
    add_textbox(s2, 2.29, 5.97, 3.0, 0.91, "ChatGPT", FONT_R, 18, True, BRAND_RGB)
    add_textbox(s2, 2.29, 6.88, 12.5, 1.78,
                "① 메모리 OFF → 검색 OFF\n    새 채팅 후 메모리·검색 모두 비활성화\n"
                "② 메모리 OFF → 검색 ON\n    새 채팅 후 검색만 활성화",
                FONT_L, 14, False, DARK)
    add_rect(s2, 16.51, 5.46, 13.23, 7.19, PPTXRGBColor(255, 248, 240))
    add_textbox(s2, 16.82, 5.97, 3.0, 0.91, "Gemini", FONT_R, 18, True, BRAND_RGB)
    add_textbox(s2, 16.82, 6.88, 12.5, 1.78,
                "① 시크릿 모드 → gemini.google.com 접속\n     로그아웃 상태에서 질문",
                FONT_L, 14, False, DARK)
    add_textbox(s2, 1.98, 13.72, 29.76, 1.91,
                "📌  각 질문은 새 채팅에서 입력하지 않고, 동일한 채팅 내에서 Q1→Q7 순서대로 연속 입력합니다.\n"
                "📌  답변은 복사해서 별도 파일에 Q번호와 함께 저장해주세요. (예: Q1_GPT검색OFF.txt)",
                FONT_L, 14, False, DARK)

    # 슬라이드 3: 질문 타이틀
    s3 = prs.slides.add_slide(blank)
    add_textbox(s3, 1.25, 2.24, 10.0, 1.45, "진단 질문 7개", FONT_R, 28, False, DARK)
    add_textbox(s3, 1.25, 3.89, 20.0, 0.94,
                f"{d['brand_name']} AI 브랜드 진단을 위한 핵심 질문",
                FONT_L, 18, False, GRAY)

    # Q1~Q7 슬라이드
    for i, q_data in enumerate(d['questions']):
        n     = i + 1
        q_txt = q_data.get('question', '')
        qtype = q_data.get('type', '')
        check = q_data.get('check_point', '')
        dlist = q_data.get('data', [])
        clean_check = check.replace('확인 포인트:', '').replace('확인 포인트 :', '').strip()

        slide = prs.slides.add_slide(blank)
        add_rect(slide, 2.11, 2.90, 1.93, 1.27, BRAND_RGB)
        add_textbox(slide, 2.11, 2.90, 1.93, 1.27,
                    f"Q{n}", FONT_R, 24, False, WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, 4.42, 2.90, 25.0, 1.27, q_txt, FONT_R, 22, False, DARK)

        add_rect(slide, 2.11, 4.60, 5.82, 1.52, BRAND_RGB)
        add_textbox(slide, 2.11, 4.60, 5.82, 1.52, qtype, FONT_L, 15, False, WHITE)
        add_textbox(slide, 8.28, 4.88, 19.94, 0.94,
                    f"확인 포인트 :  {clean_check}", FONT_L, 14, False, GRAY)

        add_textbox(slide, 2.11, 7.70, 5.31, 0.94,
                    "📊  선정 근거 데이터", FONT_L, 15, False, GRAY)

        table_data = [["출처", "데이터", "연도"]]
        for ditem in dlist:
            table_data.append([ditem.get('source', ''), ditem.get('content', ''), ditem.get('year', '')])
        while len(table_data) < 4:
            table_data.append(["", "", ""])

        tbl = slide.shapes.add_table(
            len(table_data), 3,
            PPTXCm(2.11), PPTXCm(8.82), PPTXCm(29.65), PPTXCm(3.5)).table
        tbl.columns[0].width = PPTXCm(4.0)
        tbl.columns[1].width = PPTXCm(23.65)
        tbl.columns[2].width = PPTXCm(2.0)
        for ri, row_data in enumerate(table_data):
            for ci, val in enumerate(row_data):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                if ri == 0:
                    cell.fill.fore_color.rgb = CW_LIGHT
                    clr, bold = WHITE, True
                else:
                    cell.fill.fore_color.rgb = PPTXRGBColor(255, 255, 255) if ri % 2 == 1 else PPTXRGBColor(245, 242, 237)
                    clr, bold = DARK, False
                tf = cell.text_frame
                p  = tf.paragraphs[0]
                run = p.add_run()
                run.text = val
                run.font.name  = FONT_L
                run.font.size  = PPTXPt(12)
                run.font.bold  = bold
                run.font.color.rgb = clr

        ins_text = clean_check if clean_check else f"→  이 질문에서 {d['brand_name']}의 포지셔닝을 확인하세요."
        add_rect(slide, 2.11, 15.09, 29.65, 2.11, PPTXRGBColor(234, 220, 244))
        add_textbox(slide, 2.64, 15.72, 27.51, 0.94, ins_text, FONT_L, 14, False, DARK)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────
# [4] 전략제안서 PPT 생성
# ─────────────────────────────────────────────
def create_proposal_ppt(d):
    FONT_M = "맑은 고딕"
    FONT_R = "맑은 고딕"
    FONT_L = "맑은 고딕"
    br, bg, bb = hex_to_rgb(d['brand_color'])
    BRAND_RGB   = PPTXRGBColor(br, bg, bb)
    CW_PURPLE   = PPTXRGBColor(83, 39, 168)
    CW_LIGHT    = PPTXRGBColor(124, 92, 191)
    CW_LAVENDER = PPTXRGBColor(234, 220, 244)
    DARK        = PPTXRGBColor(26, 23, 20)
    GRAY        = PPTXRGBColor(85, 85, 85)
    WHITE       = PPTXRGBColor(255, 255, 255)
    GREEN       = PPTXRGBColor(217, 242, 208)
    ORANGE_C    = PPTXRGBColor(250, 226, 213)

    ai_keys   = ['off', 'on', 'gem']

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    blank = prs.slide_layouts[6]

    def add_textbox(slide, left, top, width, height, text,
                    font_name=FONT_L, font_size=16, bold=False,
                    color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(
            PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = PPTXPt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return tb

    def add_rect(slide, left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(
            1, PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def part_slide(title_num, title_text):
        slide = prs.slides.add_slide(blank)
        add_textbox(slide, 2.82, 6.68, 6.05, 2.13,
                    f"PART {title_num}", FONT_R, 44, False, BRAND_RGB)
        add_rect(slide, 2.82, 8.81, 4.24, 0.15, BRAND_RGB)
        add_textbox(slide, 2.82, 9.22, 12.0, 1.27,
                    title_text, FONT_R, 24, False, DARK)
        return slide

    # 표지
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 33.87, 19.05, PPTXRGBColor(245, 245, 245))
    add_rect(s1, 8.56, 4.55, 0.46, 5.39, CW_PURPLE)
    add_textbox(s1, 9.38, 4.73, 12.0, 1.81, "크림웍스\nGEO 컨설팅", FONT_M, 36, False, DARK)
    add_textbox(s1, 9.38, 6.55, 7.0, 1.27, "[전략 제안서]", FONT_R, 24, False, DARK)
    add_textbox(s1, 9.38, 7.73, 14.1, 1.63,
                f"AI 검색시대, {d['brand_name']}의 ChatGPT·Gemini\n노출 현황 진단 및 GEO 개선 전략입니다.",
                FONT_R, 16, False, DARK)
    add_textbox(s1, 11.20, 17.69, 11.47, 0.86,
                "CREAMWORKS  -  AI가 좋아하는 브랜드를 만듭니다",
                FONT_M, 14, False, CW_PURPLE)

    # PART 0
    part_slide(0, "진단 질문 설계 - 근거와 방법론")
    s0 = prs.slides.add_slide(blank)
    add_textbox(s0, 2.11, 1.98, 15.0, 1.45, "0. 설계 철학 및 질문별 개요", FONT_R, 24, False, DARK)
    add_textbox(s0, 2.11, 3.50, 29.65, 1.80,
                "브랜드를 모르는 일반 소비자가 AI에게 정보를 요청하는 순간을 포착하는 것이 GEO의 출발점입니다.\n"
                "DISCOVER → CONSIDER → DECIDE 여정 순서로 질문 7개가 설계되었습니다.",
                FONT_L, 14, False, DARK)
    q_rows = [["번호", "단계", "질문 내용"]]
    for i, q_data in enumerate(d['questions']):
        q_rows.append([f"Q{i+1}", q_data.get('stage', ''), q_data.get('question', '')])
    tbl = s0.shapes.add_table(
        len(q_rows), 3,
        PPTXCm(2.11), PPTXCm(5.80), PPTXCm(29.65), PPTXCm(11.0)).table
    tbl.columns[0].width = PPTXCm(2.5)
    tbl.columns[1].width = PPTXCm(5.0)
    tbl.columns[2].width = PPTXCm(22.15)
    for ri, row_data in enumerate(q_rows):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = CW_LIGHT
                clr, bold = WHITE, True
            else:
                cell.fill.fore_color.rgb = PPTXRGBColor(255, 255, 255) if ri % 2 == 1 else PPTXRGBColor(245, 242, 237)
                clr, bold = DARK, False
            tf = cell.text_frame
            p  = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name  = FONT_L
            run.font.size  = PPTXPt(13)
            run.font.bold  = bold
            run.font.color.rgb = clr

    # PART 1: B2A 매트릭스
    part_slide(1, "AI 진단결과 - 현황분석")
    s1b = prs.slides.add_slide(blank)
    add_textbox(s1b, 2.11, 0.56, 20.0, 1.12,
                "1. B2A 매트릭스 — AI 답변 기반 언급 현황", FONT_R, 20, False, DARK)
    b2a_rows = [["질문", "GPT 검색OFF", "GPT 검색ON", "Gemini"]]
    for i, q_data in enumerate(d['questions']):
        n = i + 1
        row_data = [f"Q{n}. {q_data.get('question', '')[:28]}"]
        for key in ai_keys:
            ans     = d['answers'][key][n]
            mention = check_mention(ans, d['brand_name'])
            if ans.strip():
                row_data.append("✅ 언급됨" if mention else "❌ 미언급")
            else:
                row_data.append("—")
        b2a_rows.append(row_data)
    tbl2 = s1b.shapes.add_table(
        len(b2a_rows), 4,
        PPTXCm(2.11), PPTXCm(2.0), PPTXCm(29.65), PPTXCm(15.5)).table
    tbl2.columns[0].width = PPTXCm(12.0)
    for ci in range(1, 4):
        tbl2.columns[ci].width = PPTXCm(5.88)
    for ri, row_data in enumerate(b2a_rows):
        for ci, val in enumerate(row_data):
            cell = tbl2.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = CW_LIGHT
                clr, bold = WHITE, True
            else:
                if "✅" in val:
                    cell.fill.fore_color.rgb = GREEN
                    clr = PPTXRGBColor(21, 87, 36)
                elif "❌" in val:
                    cell.fill.fore_color.rgb = ORANGE_C
                    clr = PPTXRGBColor(180, 50, 30)
                else:
                    cell.fill.fore_color.rgb = PPTXRGBColor(255, 255, 255) if ri % 2 == 1 else PPTXRGBColor(245, 242, 237)
                    clr = GRAY
                bold = False
            tf = cell.text_frame
            p  = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name  = FONT_L
            run.font.size  = PPTXPt(13)
            run.font.bold  = bold
            run.font.color.rgb = clr

    # PART 2, 3, 4
    part_slide(2, "브랜드 GEO 기초 진단")
    part_slide(3, "크림웍스 GEO 전략 인사이트")

    for i, q_data in enumerate(d['questions']):
        n       = i + 1
        insight = d['cw_insights'][i] if i < len(d['cw_insights']) else ''
        if not insight.strip():
            continue
        si = prs.slides.add_slide(blank)
        add_rect(si, 2.11, 0.56, 29.65, 1.52, CW_LIGHT)
        add_textbox(si, 2.64, 0.79, 28.0, 1.02,
                    f"Q{n}. {q_data.get('question', '')}", FONT_R, 18, True, WHITE)
        add_rect(si, 2.11, 2.41, 29.65, 5.46, CW_LAVENDER)
        add_textbox(si, 2.64, 2.79, 27.51, 4.37, insight, FONT_L, 16, False, DARK)

    part_slide(4, "우선 실행 액션 플랜")
    sa = prs.slides.add_slide(blank)
    add_textbox(sa, 2.11, 0.56, 15.0, 1.12, "4. GEO 실행 액션 플랜", FONT_R, 20, False, DARK)
    actions = [
        ("🚨 즉시", "robots.txt, Schema Markup, 브랜드명 통일 표기"),
        ("⭐ 1개월", "FAQ 콘텐츠 7개, llms.txt, 나무위키 보강"),
        ("3개월",   "언론 PR, 채널별 콘텐츠 배포, B2A 월간 측정 시작"),
    ]
    for ai, (label, content_text) in enumerate(actions):
        top = 2.2 + ai * 4.5
        add_rect(sa, 2.11, top, 29.65, 3.81,
                 CW_LAVENDER if ai > 0 else PPTXRGBColor(26, 23, 20))
        clr = BRAND_RGB if ai == 0 else CW_PURPLE
        add_textbox(sa, 2.64, top + 0.51, 5.0, 0.94, label, FONT_R, 18, True, clr)
        add_textbox(sa, 8.0,  top + 0.51, 23.0, 2.54, content_text, FONT_L, 15, False, DARK)

    # 마지막 슬라이드
    sf = prs.slides.add_slide(blank)
    add_rect(sf, 0, 0, 33.87, 19.05, PPTXRGBColor(26, 23, 20))
    add_textbox(sf, 9.38, 8.38, 15.0, 1.45,
                "CREAMWORKS  —  AI가 좋아하는 브랜드를 만듭니다",
                FONT_M, 20, False, CW_PURPLE, PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────
# 페이지 설정
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="GEO-Report | 크림웍스",
    layout="wide",
    initial_sidebar_state="collapsed"
)

CW_PURPLE       = "#6B4EFF"
CW_PURPLE_LIGHT = "#EDE9FF"

st.markdown(f"""
<style>
  .main .block-container {{ padding-top: 1.5rem; max-width: 860px; margin: 0 auto; }}
  .brand-header {{
      padding: 24px 28px; border-radius: 14px;
      margin-bottom: 1.5rem; color: white;
  }}
  .step-label {{
      font-size: 0.72rem; font-weight: 600; color: #999;
      letter-spacing: 1px; text-transform: uppercase; margin-bottom: 2px;
  }}
  .step-title {{
      font-size: 1.1rem; font-weight: 700; color: #1a1a1a; margin-bottom: 0.8rem;
  }}
  .cw-box {{
      background: {CW_PURPLE_LIGHT};
      border-left: 4px solid {CW_PURPLE};
      padding: 14px 18px; border-radius: 0 10px 10px 0;
      margin: 8px 0 16px 0; font-size: 0.88rem; color: #3a2d7a;
  }}
  .mention-yes {{
      background: #d4edda; color: #155724;
      padding: 3px 10px; border-radius: 10px;
      font-size: 0.78rem; font-weight: 700;
  }}
  .mention-no {{
      background: #f8d7da; color: #721c24;
      padding: 3px 10px; border-radius: 10px;
      font-size: 0.78rem; font-weight: 700;
  }}
  .divider {{ border: none; border-top: 1px solid #eee; margin: 1.5rem 0; }}
  div[data-testid="stButton"] button[kind="primary"] {{
      background-color: #52B788 !important;
      border-color: #52B788 !important;
      color: white !important;
      font-weight: 600 !important;
      border-radius: 8px !important;
  }}
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────
# 세션 초기화
# ─────────────────────────────────────────────
if 'data' not in st.session_state:
    st.session_state.data = None
if 'loaded' not in st.session_state:
    st.session_state.loaded = False

# ─────────────────────────────────────────────
# 헤더
# ─────────────────────────────────────────────
brand_nm = st.session_state.data['brand_name'] if st.session_state.data else 'GEO-Report'
brand_color = st.session_state.data['brand_color'] if st.session_state.data else '#4A90D9'

st.markdown(f"""
<div class="brand-header" style="background: linear-gradient(135deg, #0f0f0f 60%, {brand_color}99);">
  <div style="font-size:0.78rem;color:#888;letter-spacing:2px;font-weight:600;">CREAMWORKS  ·  GEO-Report  |  앱 2</div>
  <div style="font-size:1.7rem;font-weight:800;margin:6px 0 4px;">{brand_nm} 보고서 생성</div>
  <div style="font-size:0.85rem;color:#bbb;">GEO-Scan Excel 업로드 → 질문지·전략제안서 Word·PPT 자동 생성</div>
</div>
""", unsafe_allow_html=True)

st.markdown("<hr class='divider'>", unsafe_allow_html=True)

# ─────────────────────────────────────────────
# STEP 1: Excel 업로드
# ─────────────────────────────────────────────
if not st.session_state.loaded:
    st.markdown('<div class="step-label">STEP 1</div>', unsafe_allow_html=True)
    st.markdown('<div class="step-title">GEO-Scan Excel 업로드</div>', unsafe_allow_html=True)
    st.markdown('<div class="cw-box">💜 GEO-Scan 앱(앱1)에서 저장한 Excel 파일을 업로드하세요. 브랜드 정보·질문·답변·분석 결과가 자동으로 불러와집니다.</div>', unsafe_allow_html=True)

    uploaded = st.file_uploader(
        "Excel 파일 업로드 (.xlsx)",
        type=["xlsx"],
        key="excel_upload"
    )

    if uploaded:
        try:
            with st.spinner("Excel 파일 읽는 중..."):
                d = load_from_excel(uploaded)
            st.session_state.data   = d
            st.session_state.loaded = True
            st.success(f"✅ 업로드 완료! 브랜드: **{d['brand_name']}** | 질문 **{len(d['questions'])}개** | 분석결과 {'있음' if d['analysis_result'] else '없음'}")
            st.rerun()
        except Exception as e:
            st.error(f"파일 읽기 오류: {e}")

# ─────────────────────────────────────────────
# STEP 2: 보고서 생성
# ─────────────────────────────────────────────
else:
    d = st.session_state.data

    st.markdown('<div class="step-label">STEP 2</div>', unsafe_allow_html=True)
    st.markdown('<div class="step-title">보고서 생성</div>', unsafe_allow_html=True)

    # 브랜드 정보 요약
    with st.expander("📋 불러온 브랜드 정보 확인", expanded=False):
        col1, col2 = st.columns(2)
        with col1:
            st.markdown(f"**브랜드명:** {d['brand_name']}")
            st.markdown(f"**카테고리:** {d['brand_category']}")
            st.markdown(f"**브랜드 컬러:** {d['brand_color']}")
            st.markdown(f"**경쟁 브랜드:** {d['brand_competitors']}")
        with col2:
            st.markdown(f"**핵심 USP:** {d['brand_usp'][:80]}{'...' if len(d['brand_usp']) > 80 else ''}")
            st.markdown(f"**질문 수:** {len(d['questions'])}개")
            ai_keys = ['off', 'on', 'gem']
            total_ans = sum(1 for key in ai_keys for n in range(1, 8) if d['answers'][key][n].strip())
            st.markdown(f"**수집된 답변:** {total_ans}/21개")
            st.markdown(f"**분석 결과:** {'있음 ✅' if d['analysis_result'] else '없음 ⚠️'}")

    st.markdown('<div class="cw-box">💜 브랜드 컬러 <b>' + d['brand_color'] + '</b> 자동 적용 &nbsp;|&nbsp; 크림웍스 퍼플 <b>#7030A0</b> &nbsp;|&nbsp; 폰트 <b>맑은 고딕</b></div>', unsafe_allow_html=True)

    # B2A 현황 미리보기
    st.markdown("#### 📊 B2A 언급 현황")
    ai_keys       = ['off', 'on', 'gem']
    ai_labels_s   = ['GPT OFF', 'GPT ON', 'Gemini']
    header = st.columns([4, 1, 1, 1])
    header[0].markdown("**질문**")
    for j, lbl in enumerate(ai_labels_s):
        header[j+1].markdown(f"**{lbl}**")
    st.markdown("<hr style='border:none;border-top:1px solid #ddd;margin:4px 0'>", unsafe_allow_html=True)
    for i, q_data in enumerate(d['questions']):
        n = i + 1
        row = st.columns([4, 1, 1, 1])
        q_txt = q_data.get('question', '')
        row[0].markdown(f"**Q{n}.** {q_txt[:40]}{'...' if len(q_txt) > 40 else ''}")
        for j, key in enumerate(ai_keys):
            ans = d['answers'][key][n]
            mention = check_mention(ans, d['brand_name'])
            with row[j+1]:
                if ans.strip():
                    if mention:
                        st.markdown('<span class="mention-yes">✓ 언급</span>', unsafe_allow_html=True)
                    else:
                        st.markdown('<span class="mention-no">✗ 미언급</span>', unsafe_allow_html=True)
                else:
                    st.caption("—")

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    # ── 질문지 생성 ──
    st.markdown("#### 📄 질문지")
    col_qw, col_qp = st.columns(2)

    with col_qw:
        st.markdown("**질문지 Word** — 내부 검토용")
        if st.button("📄 질문지 Word 생성", use_container_width=True, key="btn_qw"):
            with st.spinner("생성 중..."):
                buf = create_question_word(d)
            fname = f"{d['brand_name']}_GEO_AI진단질문지_{datetime.now().strftime('%Y%m%d')}.docx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True, key="dl_qw"
            )

    with col_qp:
        st.markdown("**질문지 PPT** — 광고주 제안용")
        if st.button("📊 질문지 PPT 생성", use_container_width=True, key="btn_qp"):
            with st.spinner("생성 중..."):
                buf = create_question_ppt(d)
            fname = f"{d['brand_name']}_GEO_AI진단질문지_{datetime.now().strftime('%Y%m%d')}.pptx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True, key="dl_qp"
            )

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    # ── 전략제안서 생성 ──
    st.markdown("#### 📑 전략제안서")
    col_pw, col_pp = st.columns(2)

    with col_pw:
        st.markdown("**전략제안서 Word** — 내부 검토용")
        if st.button("📄 전략제안서 Word 생성", use_container_width=True, key="btn_pw"):
            with st.spinner("생성 중..."):
                buf = create_proposal_word(d)
            fname = f"{d['brand_name']}_GEO_전략제안서_{datetime.now().strftime('%Y%m%d')}.docx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True, key="dl_pw"
            )

    with col_pp:
        st.markdown("**전략제안서 PPT** — 광고주 제안용")
        if st.button("📊 전략제안서 PPT 생성", use_container_width=True, key="btn_pp"):
            with st.spinner("생성 중..."):
                buf = create_proposal_ppt(d)
            fname = f"{d['brand_name']}_GEO_전략제안서_{datetime.now().strftime('%Y%m%d')}.pptx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True, key="dl_pp"
            )

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    # ── 다른 파일 열기 ──
    col_reset, _ = st.columns([1, 3])
    with col_reset:
        if st.button("🔄 다른 Excel 파일 열기", use_container_width=True):
            st.session_state.data   = None
            st.session_state.loaded = False
            st.rerun()
