import streamlit as st
import anthropic
from docx import Document
from docx.shared import RGBColor, Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PPTXPt, Emu, Cm as PPTXCm
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.text import PP_ALIGN
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import io
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side


def brand_rgb():
    return hex_to_rgb(st.session_state.brand_color)


# ─────────────────────────────────────────────
# Excel 불러오기 함수
# ─────────────────────────────────────────────
def load_from_excel(uploaded_file):
    """
    Excel 파일을 읽어 session_state에 복원.
    반환값: 'answer' (답변수집) 또는 'full' (전체데이터)
    """
    wb = openpyxl.load_workbook(uploaded_file)
    sheets = wb.sheetnames

    # ── 브랜드 정보 복원 ──
    if '브랜드 정보' in sheets:
        ws = wb['브랜드 정보']
        info_map = {
            '브랜드명':   'brand_name',
            '카테고리':   'brand_category',
            '핵심 USP':  'brand_usp',
            '주요 타겟':  'brand_target',
            '경쟁 브랜드': 'brand_competitors',
            '부정 이미지': 'brand_negative',
            '강조 포인트': 'brand_focus',
        }
        for row in ws.iter_rows(min_row=2, values_only=True):
            label, value = row[0], row[1]
            if label in info_map and value:
                st.session_state[info_map[label]] = str(value)

    # ── 질문 + 답변 복원 ──
    if 'AI 답변 수집' in sheets:
        ws2 = wb['AI 답변 수집']
        questions = []
        answers = {
            'off': {i: '' for i in range(1, 8)},
            'on':  {i: '' for i in range(1, 8)},
            'gem': {i: '' for i in range(1, 8)},
        }
        for row in ws2.iter_rows(min_row=2, values_only=True):
            if not row[0] or not str(row[0]).startswith('Q'):
                continue
            n = int(str(row[0]).replace('Q', ''))
            # 컬럼: 번호(0), 질문(1), 유형(2), 단계(3), 확인포인트(4)
            #        GPT OFF(5), GPT ON(6), Gemini(7), 언급(8,9,10)
            q_dict = {
                'question':    str(row[1]) if row[1] else '',
                'type':        str(row[2]) if row[2] else '',
                'stage':       str(row[3]) if row[3] else '',
                'check_point': str(row[4]) if len(row) > 4 and row[4] else '',
                'data':        [],
            }
            questions.append(q_dict)
            # 새 형식 (확인포인트 컬럼 추가됨): 5,6,7
            # 구 형식 (확인포인트 없음): 4,5,6
            if len(row) >= 8 and row[5] is not None:
                answers['off'][n] = str(row[5]) if row[5] else ''
                answers['on'][n]  = str(row[6]) if row[6] else ''
                answers['gem'][n] = str(row[7]) if row[7] else ''
            else:
                answers['off'][n] = str(row[4]) if len(row) > 4 and row[4] else ''
                answers['on'][n]  = str(row[5]) if len(row) > 5 and row[5] else ''
                answers['gem'][n] = str(row[6]) if len(row) > 6 and row[6] else ''

        st.session_state.questions = questions
        st.session_state.answers   = answers
        st.session_state.questions_confirmed = True

    # ── Claude 분석 결과 복원 ──
    has_analysis = False
    if 'Claude 분석' in sheets:
        ws4 = wb['Claude 분석']
        for row in ws4.iter_rows(min_row=2, values_only=True):
            if row[0]:
                st.session_state.analysis_result = str(row[0])
                has_analysis = True
                break

    # ── 파일 타입 판별 ──
    if has_analysis and st.session_state.get('analysis_result', '').strip():
        # 전체데이터: STEP 5로
        st.session_state.overall_diagnosis  = st.session_state.analysis_result[:300]
        st.session_state.cw_insights        = [''] * 7
        st.session_state.priority_actions   = ''
        st.session_state.step = 5
        return 'full'
    else:
        # 답변수집: STEP 4로
        st.session_state.analysis_result    = ''
        st.session_state.overall_diagnosis  = ''
        st.session_state.cw_insights        = [''] * 7
        st.session_state.priority_actions   = ''
        st.session_state.step = 4
        return 'answer'



# ─────────────────────────────────────────────
# Excel 답변 저장 함수
# ─────────────────────────────────────────────
def create_answer_excel():
    br, bg, bb = brand_rgb()
    BRAND_HEX = st.session_state.brand_color.replace('#', 'FF')
    CW_HEX    = 'FFEADCF4'
    DARK_HEX  = 'FF1A1A1A'
    GRAY_HEX  = 'FFF7F7F7'
    GREEN_HEX = 'FFD9F2D0'
    RED_HEX   = 'FFFAE2D5'
    WHITE_HEX = 'FFFFFFFF'
    GRAY2_HEX = 'FFFAFAFA'
    # 헤더용: 브랜드 컬러 연하게 (흰 글씨 대신 어두운 글씨 사용)
    HEADER_HEX = f'FF{br:02X}{bg:02X}{bb:02X}'

    wb = openpyxl.Workbook()

    # ── 시트 1: 브랜드 정보 ──
    ws_info = wb.active
    ws_info.title = "브랜드 정보"

    header_fill = PatternFill("solid", fgColor=HEADER_HEX)
    brand_fill  = PatternFill("solid", fgColor=BRAND_HEX)
    cw_fill     = PatternFill("solid", fgColor=CW_HEX)
    gray_fill   = PatternFill("solid", fgColor=GRAY_HEX)
    white_fill  = PatternFill("solid", fgColor=WHITE_HEX)
    green_fill  = PatternFill("solid", fgColor=GREEN_HEX)
    red_fill    = PatternFill("solid", fgColor=RED_HEX)
    gray2_fill  = PatternFill("solid", fgColor=GRAY2_HEX)

    thin = Side(style='thin', color='DDDDDD')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def cell_style(ws, row, col, value, fill=None, font_color='FF000000',
                   bold=False, align='left', font_size=10):
        c = ws.cell(row=row, column=col, value=value)
        if fill:
            c.fill = fill
        c.font = Font(name='맑은 고딕', size=font_size,
                      bold=bold, color=font_color)
        c.alignment = Alignment(horizontal=align, vertical='center',
                                wrap_text=True)
        c.border = border
        return c

    # 브랜드 정보 헤더
    ws_info.merge_cells('A1:B1')
    c = ws_info['A1']
    c.value = f"{st.session_state.brand_name}  GEO 진단 — 브랜드 정보"
    c.fill = header_fill
    c.font = Font(name='맑은 고딕', size=13, bold=True, color='FFFFFFFF')
    c.alignment = Alignment(horizontal='left', vertical='center')
    ws_info.row_dimensions[1].height = 30

    info_data = [
        ("브랜드명",    st.session_state.brand_name),
        ("카테고리",    st.session_state.brand_category),
        ("핵심 USP",   st.session_state.brand_usp),
        ("주요 타겟",  st.session_state.brand_target),
        ("경쟁 브랜드", st.session_state.brand_competitors),
        ("부정 이미지", st.session_state.brand_negative),
        ("강조 포인트", st.session_state.brand_focus),
        ("진단일",     datetime.now().strftime('%Y.%m.%d')),
    ]
    for i, (label, value) in enumerate(info_data):
        row = i + 2
        cell_style(ws_info, row, 1, label, brand_fill,
                   f'FF{br:02X}{bg:02X}{bb:02X}', True, 'left', 10)
        cell_style(ws_info, row, 2, value, gray_fill if i%2==0 else white_fill,
                   'FF333333', False, 'left', 10)
        ws_info.row_dimensions[row].height = 40

    ws_info.column_dimensions['A'].width = 18
    ws_info.column_dimensions['B'].width = 60

    # ── 시트 2: 질문 + 답변 Raw Data ──
    ws_ans = wb.create_sheet("AI 답변 수집")

    ai_keys   = ['off', 'on', 'gem']
    ai_labels = ['ChatGPT 검색OFF', 'ChatGPT 검색ON', 'Gemini']

    # 헤더 행
    headers = ["번호", "질문", "유형", "단계", "확인포인트",
               "ChatGPT 검색OFF", "ChatGPT 검색ON", "Gemini",
               "GPT OFF 언급", "GPT ON 언급", "Gemini 언급"]
    for ci, h in enumerate(headers):
        cell_style(ws_ans, 1, ci+1, h, header_fill,
                   'FFFFFFFF', True, 'center', 10)
    ws_ans.row_dimensions[1].height = 25

    # 데이터 행
    for i, q_data in enumerate(st.session_state.questions):
        n     = i + 1
        row   = i + 2
        bg_fill = gray_fill if i%2==0 else white_fill

        cell_style(ws_ans, row, 1, f"Q{n}", bg_fill, 'FF333333', True, 'center', 10)
        cell_style(ws_ans, row, 2, q_data.get('question',''), bg_fill, 'FF333333', False, 'left', 10)
        cell_style(ws_ans, row, 3, q_data.get('type',''), bg_fill, 'FF555555', False, 'left', 9)
        cell_style(ws_ans, row, 4, q_data.get('stage',''), bg_fill, 'FF555555', False, 'center', 9)
        cell_style(ws_ans, row, 5, q_data.get('check_point',''), bg_fill, 'FF555555', False, 'left', 9)

        for j, key in enumerate(ai_keys):
            ans = st.session_state.answers[key][n]
            cell_style(ws_ans, row, 6+j, ans, bg_fill, 'FF333333', False, 'left', 9)

            # 언급 여부
            mention = check_mention(ans, st.session_state.brand_name) if ans.strip() else None
            if mention is True:
                cell_style(ws_ans, row, 9+j, "✅ 언급됨", green_fill, 'FF155724', True, 'center', 9)
            elif mention is False:
                cell_style(ws_ans, row, 9+j, "❌ 미언급", red_fill, 'FFB43216', True, 'center', 9)
            else:
                cell_style(ws_ans, row, 9+j, "—", bg_fill, 'FF999999', False, 'center', 9)

        ws_ans.row_dimensions[row].height = 80

    col_widths = [8, 40, 25, 12, 35, 50, 50, 50, 12, 12, 12]
    for ci, w in enumerate(col_widths):
        ws_ans.column_dimensions[openpyxl.utils.get_column_letter(ci+1)].width = w

    # ── 시트 3: B2A 분석 결과 ──
    ws_b2a = wb.create_sheet("B2A 분석 결과")

    ws_b2a.merge_cells('A1:D1')
    c = ws_b2a['A1']
    c.value = f"B2A 매트릭스 — {st.session_state.brand_name}  |  수집일: {datetime.now().strftime('%Y.%m.%d')}"
    c.fill = header_fill
    c.font = Font(name='맑은 고딕', size=12, bold=True, color='FFFFFFFF')
    c.alignment = Alignment(horizontal='left', vertical='center')
    ws_b2a.row_dimensions[1].height = 28

    b2a_headers = ["질문", "GPT 검색OFF", "GPT 검색ON", "Gemini"]
    for ci, h in enumerate(b2a_headers):
        cell_style(ws_b2a, 2, ci+1, h, cw_fill, 'FF5327A8', True, 'center', 10)
    ws_b2a.row_dimensions[2].height = 22

    for i, q_data in enumerate(st.session_state.questions):
        n   = i + 1
        row = i + 3
        bg_fill = gray2_fill if i%2==0 else white_fill

        cell_style(ws_b2a, row, 1,
                   f"Q{n}. {q_data.get('question','')}",
                   bg_fill, 'FF333333', False, 'left', 9)

        for j, key in enumerate(ai_keys):
            ans     = st.session_state.answers[key][n]
            mention = check_mention(ans, st.session_state.brand_name) if ans.strip() else None
            if mention is True:
                cell_style(ws_b2a, row, 2+j, "✅ 언급됨", green_fill, 'FF155724', True, 'center', 9)
            elif mention is False:
                cell_style(ws_b2a, row, 2+j, "❌ 미언급", red_fill, 'FFB43216', True, 'center', 9)
            else:
                cell_style(ws_b2a, row, 2+j, "—", bg_fill, 'FF999999', False, 'center', 9)

        ws_b2a.row_dimensions[row].height = 30

    ws_b2a.column_dimensions['A'].width = 50
    for col in ['B','C','D']:
        ws_b2a.column_dimensions[col].width = 18

    # ── 시트 4: Claude 분석 결과 ──
    ws_cl = wb.create_sheet("Claude 분석")
    ws_cl.merge_cells('A1:B1')
    c = ws_cl['A1']
    c.value = "Claude B2A 분석 결과"
    c.fill = header_fill
    c.font = Font(name='맑은 고딕', size=12, bold=True, color='FFFFFFFF')
    c.alignment = Alignment(horizontal='left', vertical='center')
    ws_cl.row_dimensions[1].height = 28

    analysis_text = st.session_state.get('analysis_result', '')
    c2 = ws_cl.cell(row=2, column=1, value=analysis_text)
    c2.font = Font(name='맑은 고딕', size=9)
    c2.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
    c2.fill = PatternFill("solid", fgColor=CW_HEX)
    ws_cl.row_dimensions[2].height = 300
    ws_cl.column_dimensions['A'].width = 100

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf

# ─────────────────────────────────────────────
# [1] 질문지 Word 생성
# ─────────────────────────────────────────────
def create_question_word():
    from docx.oxml.ns import qn as docx_qn
    from docx.oxml import OxmlElement
    from docx.enum.table import WD_TABLE_ALIGNMENT

    FONT    = "페이퍼로지 3 Light"
    br, bg, bb = brand_rgb()
    # 크림웍스 퍼플
    cr, cg, cb = (112, 48, 160)
    BRAND_HEX  = st.session_state.brand_color.replace('#','')

    # 브랜드 컬러 30% 연하게 (흰색과 혼합)
    def lighten(r_, g_, b_, factor=0.5):
        return (
            int(r_ + (255 - r_) * factor),
            int(g_ + (255 - g_) * factor),
            int(b_ + (255 - b_) * factor),
        )
    lr, lg, lb = lighten(br, bg, bb, 0.5)
    BRAND_LIGHT_HEX = f'{lr:02X}{lg:02X}{lb:02X}'  # 유형 박스 배경

    CW_HEX     = 'EADCF4'   # 연보라 (헤더/강조)
    GRAY_HEX   = 'F7F7F7'
    WHITE_HEX  = 'FFFFFF'
    GRAY2_HEX  = 'FAFAFA'
    NOTICE_HEX = 'EADCF4'   # 주의사항도 연보라로 통일
    DARK_HEX   = '1A1A14'

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Cm(1.8)
        section.bottom_margin = Cm(1.8)
        section.left_margin   = Cm(2.2)
        section.right_margin  = Cm(2.2)

    def set_bg(cell, hex_color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(docx_qn('w:val'), 'clear')
        shd.set(docx_qn('w:color'), 'auto')
        shd.set(docx_qn('w:fill'), hex_color)
        tcPr.append(shd)

    def r(para, text, size=10, bold=False, color=None, font=FONT, italic=False):
        run = para.add_run(text)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.bold        = bold
        run.font.italic = italic
        if color:
            try:
                run.font.color.rgb = RGBColor(int(color[0]), int(color[1]), int(color[2]))
            except Exception:
                pass
        return run

    def set_cell_margins(cell, top=80, bottom=80, left=120, right=120):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcMar = OxmlElement('w:tcMar')
        for side, val in [('top',top),('bottom',bottom),('left',left),('right',right)]:
            node = OxmlElement(f'w:{side}')
            node.set(docx_qn('w:w'), str(val))
            node.set(docx_qn('w:type'), 'dxa')
            tcMar.append(node)
        tcPr.append(tcMar)

    def add_border_line(doc, color_hex, thickness=6):
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bot = OxmlElement('w:bottom')
        bot.set(docx_qn('w:val'), 'single')
        bot.set(docx_qn('w:sz'), str(thickness))
        bot.set(docx_qn('w:space'), '1')
        bot.set(docx_qn('w:color'), color_hex)
        pBdr.append(bot)
        pPr.append(pBdr)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(4)
        return p

    # ── 표지 ──
    p_header = doc.add_paragraph()
    p_header.paragraph_format.space_before = Pt(0)
    p_header.paragraph_format.space_after  = Pt(2)
    r(p_header, "CREAMWORKS  |  GEO 컨설팅 제안서  |  Confidential",
      size=9, color=(85,85,85))

    add_border_line(doc, BRAND_HEX, thickness=6)

    p_brand = doc.add_paragraph()
    p_brand.paragraph_format.space_before = Pt(40)
    p_brand.paragraph_format.space_after  = Pt(6)
    r(p_brand, st.session_state.brand_name, size=36, color=(26,23,20))

    p_sub = doc.add_paragraph()
    p_sub.paragraph_format.space_before = Pt(0)
    p_sub.paragraph_format.space_after  = Pt(6)
    r(p_sub, "AI 진단 질문지", size=18, color=(cr,cg,cb))

    p_date = doc.add_paragraph()
    p_date.paragraph_format.space_before = Pt(0)
    p_date.paragraph_format.space_after  = Pt(40)
    r(p_date, f"Presented by CREAMWORKS  ·  {datetime.now().strftime('%Y.%m')}",
      size=10, color=(85,85,85))

    doc.add_page_break()

    # ── 실행 전 필수 세팅 ──
    p_h = doc.add_paragraph()
    p_h.paragraph_format.space_after = Pt(8)
    r(p_h, "실행 전 필수 세팅", size=14, bold=True, color=(26,23,20))

    setup_table = doc.add_table(rows=2, cols=2)
    setup_table.style = 'Table Grid'
    setup_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    # 열 너비 설정
    from docx.shared import Cm as DocxCm
    for row in setup_table.rows:
        row.cells[0].width = DocxCm(3.0)
        row.cells[1].width = DocxCm(13.0)

    setup_data = [
        ("ChatGPT",
         "① 메모리 OFF → 검색 OFF  (새 채팅 시작 후 메모리·검색 모두 비활성화)\n② 메모리 OFF → 검색 ON  (새 채팅 시작 후 검색만 활성화)"),
        ("Gemini",
         "③ 시크릿 모드 → gemini.google.com 접속  (로그아웃 상태에서 질문)"),
    ]
    for ri, (label, content) in enumerate(setup_data):
        set_bg(setup_table.rows[ri].cells[0], DARK_HEX)
        set_bg(setup_table.rows[ri].cells[1], GRAY_HEX)
        set_cell_margins(setup_table.rows[ri].cells[0])
        set_cell_margins(setup_table.rows[ri].cells[1])
        rp = setup_table.rows[ri].cells[0].paragraphs[0]
        r(rp, label, size=10, bold=True, color=(br,bg,bb))
        cp = setup_table.rows[ri].cells[1].paragraphs[0]
        r(cp, content, size=9.5, color=(26,23,20))

    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # 주의사항 박스
    notice_t = doc.add_table(rows=1, cols=1)
    notice_t.style = 'Table Grid'
    set_bg(notice_t.rows[0].cells[0], NOTICE_HEX)
    set_cell_margins(notice_t.rows[0].cells[0])
    np_ = notice_t.rows[0].cells[0].paragraphs[0]
    r(np_, "📌  각 질문은 새 채팅에서 입력하지 않고, 동일한 채팅 내에서 Q1→Q7 순서대로 연속 입력합니다.\n"
           "📌  답변은 복사해서 별도 파일에 Q번호와 함께 저장해주세요. (예: Q1_GPT검색OFF.txt)",
      size=9.5, color=(26,23,20))

    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    p_h2 = doc.add_paragraph()
    p_h2.paragraph_format.space_after = Pt(8)
    r(p_h2, "진단 질문 7개", size=14, bold=True, color=(26,23,20))

    # ── Q1~Q7 ──
    for i, q_data in enumerate(st.session_state.questions):
        n      = i + 1
        q_txt  = q_data.get('question','')
        qtype  = q_data.get('type','')
        check  = q_data.get('check_point','')
        dlist  = q_data.get('data',[])

        # Q 헤더 행 (브랜드컬러 배경)
        q_hdr = doc.add_table(rows=1, cols=1)
        q_hdr.style = 'Table Grid'
        set_bg(q_hdr.rows[0].cells[0], BRAND_HEX)
        set_cell_margins(q_hdr.rows[0].cells[0])
        qp = q_hdr.rows[0].cells[0].paragraphs[0]
        r(qp, f"Q{n}.", size=12, bold=True, color=(255,255,255))
        r(qp, f"  {q_txt}", size=12, bold=True, color=(255,255,255))

        # 유형 + 확인포인트 행
        type_t = doc.add_table(rows=1, cols=2)
        type_t.style = 'Table Grid'
        set_bg(type_t.rows[0].cells[0], BRAND_LIGHT_HEX)
        set_bg(type_t.rows[0].cells[1], GRAY_HEX)
        set_cell_margins(type_t.rows[0].cells[0])
        set_cell_margins(type_t.rows[0].cells[1])
        tp = type_t.rows[0].cells[0].paragraphs[0]
        r(tp, "유형\n", size=9, color=(85,85,85))
        r(tp, qtype, size=9, color=(85,85,85))
        cp_ = type_t.rows[0].cells[1].paragraphs[0]
        r(cp_, f"확인 포인트 : {check}", size=9, color=(85,85,85))

        # 데이터 표 헤더
        data_rows = 1 + len(dlist)
        dt = doc.add_table(rows=data_rows, cols=3)
        dt.style = 'Table Grid'
        for ci, h_txt in enumerate(["출처","데이터","연도"]):
            set_bg(dt.rows[0].cells[ci], CW_HEX)
            set_cell_margins(dt.rows[0].cells[ci])
            r(dt.rows[0].cells[ci].paragraphs[0],
              h_txt, size=9, bold=True, color=(85,85,85))

        # 데이터 행
        for j, d in enumerate(dlist):
            row_bg = WHITE_HEX if j%2==0 else GRAY2_HEX
            for ci, val in enumerate([
                d.get('source',''), d.get('content',''), d.get('year','')
            ]):
                set_bg(dt.rows[j+1].cells[ci], row_bg)
                set_cell_margins(dt.rows[j+1].cells[ci])
                r(dt.rows[j+1].cells[ci].paragraphs[0],
                  val, size=9, color=(26,23,20))

        # 인사이트 박스 (연보라 배경) — check_point 내용 사용
        ins_t = doc.add_table(rows=1, cols=1)
        ins_t.style = 'Table Grid'
        set_bg(ins_t.rows[0].cells[0], CW_HEX)
        set_cell_margins(ins_t.rows[0].cells[0])
        ip = ins_t.rows[0].cells[0].paragraphs[0]
        r(ip, "→  ", size=9, bold=True, color=(cr,cg,cb))
        # check_point가 있으면 사용, 없으면 기본 문구
        insight_text = check if check.strip() else f"이 질문에서 {st.session_state.brand_name}이(가) 어떻게 언급되는지 + 경쟁사 대비 포지션을 확인하세요."
        r(ip, insight_text, size=9, color=(26,23,20))

        sp = doc.add_paragraph()
        sp.paragraph_format.space_after = Pt(8)

    # ── 질문 선정 근거 요약 ──
    doc.add_page_break()
    p_sum = doc.add_paragraph()
    p_sum.paragraph_format.space_after = Pt(4)
    r(p_sum, "질문 선정 근거 요약", size=14, bold=True, color=(26,23,20))

    p_sum_desc = doc.add_paragraph()
    p_sum_desc.paragraph_format.space_after = Pt(8)
    r(p_sum_desc, "이 7개 질문은 다음 자료를 교차 분석해 도출했습니다.", size=9.5, color=(85,85,85))

    sum_t = doc.add_table(rows=1, cols=3)
    sum_t.style = 'Table Grid'
    for ci, h_txt in enumerate(["자료명", "발행처", "연도"]):
        set_bg(sum_t.rows[0].cells[ci], CW_HEX)
        set_cell_margins(sum_t.rows[0].cells[ci])
        r(sum_t.rows[0].cells[ci].paragraphs[0],
          h_txt, size=9, bold=True, color=(85,85,85))

    # 각 질문의 데이터 소스 취합 (자료명: content, 발행처: source, 연도: year)
    all_sources = {}  # key: source → (content, year)
    for q_data in st.session_state.questions:
        for d in q_data.get('data', []):
            src = d.get('source', '').strip()
            if src and src not in all_sources:
                all_sources[src] = (d.get('content', ''), d.get('year', ''))

    for idx, (src, (content_val, yr)) in enumerate(all_sources.items()):
        row_bg = WHITE_HEX if idx % 2 == 0 else GRAY2_HEX
        new_row = sum_t.add_row()
        for ci in range(3):
            set_bg(new_row.cells[ci], row_bg)
            set_cell_margins(new_row.cells[ci])
        # 자료명: content 내용 요약 (없으면 source 사용)
        data_name = content_val[:40] + "…" if len(content_val) > 40 else content_val
        r(new_row.cells[0].paragraphs[0], data_name or src, size=9, color=(26,23,20))
        r(new_row.cells[1].paragraphs[0], src, size=9, color=(26,23,20))
        r(new_row.cells[2].paragraphs[0], yr, size=9, color=(26,23,20))

    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # 푸터
    fp = doc.add_paragraph()
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r(fp, "CREAMWORKS  —  AI가 좋아하는 브랜드를 만듭니다", size=9, color=(150,150,150))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def create_question_ppt():
    FONT_M  = "페이퍼로지 5 Medium"
    FONT_R  = "페이퍼로지 4 Regular"
    FONT_L  = "페이퍼로지 3 Light"
    br, bg, bb = brand_rgb()
    BRAND_RGB = PPTXRGBColor(br, bg, bb)
    CW_PURPLE = PPTXRGBColor(83, 39, 168)   # #5327A8
    CW_LIGHT  = PPTXRGBColor(124, 92, 191)  # #7C5CBF
    DARK      = PPTXRGBColor(26, 23, 20)
    GRAY      = PPTXRGBColor(85, 85, 85)
    WHITE     = PPTXRGBColor(255, 255, 255)
    YELLOW    = PPTXRGBColor(255, 204, 102)  # #FFCC66

    W = Inches(13.33)
    H = Inches(7.50)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]

    def add_textbox(slide, left, top, width, height, text, font_name=FONT_L,
                    font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(
            PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = PPTXPt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def add_rect(slide, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def add_table_ppt(slide, left, top, width, rows, cols, data, header_bg, row_bgs):
        table = slide.shapes.add_table(rows, cols,
            PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(3.5)).table
        col_widths = [PPTXCm(3.0), PPTXCm(width-5.0), PPTXCm(2.0)]
        for ci in range(min(cols, len(col_widths))):
            table.columns[ci].width = col_widths[ci]
        for ri, row_data in enumerate(data):
            bg = header_bg if ri == 0 else row_bgs[ri % 2]
            for ci, val in enumerate(row_data):
                cell = table.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = val
                run.font.name = FONT_L
                run.font.size = PPTXPt(12)
                if ri == 0:
                    run.font.bold  = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.color.rgb = DARK
        return table

    # ── 슬라이드 1: 표지 (이미지형) ──
    slide1 = prs.slides.add_slide(blank_layout)
    cover_rect = add_rect(slide1, 0, 0, 33.87, 19.05,
                          PPTXRGBColor(245, 245, 245))

    # 세로 강조선
    add_rect(slide1, 8.56, 4.55, 0.46, 5.39, CW_PURPLE)

    # 제목
    add_textbox(slide1, 9.38, 4.73, 12.0, 1.81,
                "크림웍스\nGEO 컨설팅", FONT_M, 36, False, DARK)
    add_textbox(slide1, 9.38, 6.55, 7.0, 1.27,
                "[2. AI 진단 설문지]", FONT_R, 24, False, DARK)
    add_textbox(slide1, 9.38, 7.73, 14.1, 1.63,
                f"AI 검색시대, 브랜드가 ChatGPT·Gemini에서 발견되고\n추천되기 위해 체계적인 설문지를 만드는 프로세스입니다",
                FONT_R, 16, False, DARK)
    add_textbox(slide1, 11.20, 17.69, 11.47, 0.86,
                "CREAMWORKS  - AI가 좋아하는 브랜드를 만듭니다",
                FONT_M, 14, False, CW_PURPLE)

    # ── 슬라이드 2: 실행 세팅 ──
    slide2 = prs.slides.add_slide(blank_layout)
    add_textbox(slide2, 1.25, 2.24, 7.32, 1.45,
                "실행 전 필수 세팅", FONT_R, 28, False, DARK)

    # 구분선
    line = slide2.shapes.add_shape(1,
        PPTXCm(1.25), PPTXCm(3.71), PPTXCm(31.37), PPTXCm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_RGB
    line.line.fill.background()

    # ChatGPT 박스
    add_rect(slide2, 1.98, 5.46, 13.23, 7.19, PPTXRGBColor(255, 248, 240))
    add_textbox(slide2, 2.29, 5.97, 3.0, 0.91, "ChatGPT", FONT_R, 18, True, BRAND_RGB)
    add_textbox(slide2, 2.29, 6.88, 12.5, 1.78,
                "① 메모리 OFF → 검색 OFF\n    새 채팅 시작 후 메모리·검색 모두 비활성화\n② 메모리 OFF → 검색 ON\n    새 채팅 시작 후 검색만 활성화",
                FONT_L, 14, False, DARK)

    # Gemini 박스
    add_rect(slide2, 16.51, 5.46, 13.23, 7.19, PPTXRGBColor(255, 248, 240))
    add_textbox(slide2, 16.82, 5.97, 3.0, 0.91, "Gemini", FONT_R, 18, True, BRAND_RGB)
    add_textbox(slide2, 16.82, 6.88, 12.5, 1.78,
                "① 시크릿 모드 → gemini.google.com 접속\n     로그아웃 상태에서 질문",
                FONT_L, 14, False, DARK)

    # 주의사항
    add_textbox(slide2, 1.98, 13.72, 29.76, 1.91,
                "📌  각 질문은 새 채팅에서 입력하지 않고, 동일한 채팅 내에서 Q1→Q7 순서대로 연속 입력합니다.\n"
                "📌  답변은 복사해서 별도 파일에 Q번호와 함께 저장해주세요. (예: Q1_GPT검색OFF.txt)",
                FONT_L, 14, False, DARK)

    # ── 슬라이드 3: 질문 타이틀 ──
    slide3 = prs.slides.add_slide(blank_layout)
    add_textbox(slide3, 1.25, 2.24, 10.0, 1.45,
                "진단 질문 7개", FONT_R, 28, False, DARK)
    add_textbox(slide3, 1.25, 3.89, 15.0, 0.94,
                f"{st.session_state.brand_name} AI 브랜드 진단을 위한 핵심 질문",
                FONT_L, 18, False, GRAY)

    # ── Q 슬라이드 Q1~Q7 ──
    for i, q_data in enumerate(st.session_state.questions):
        n     = i + 1
        q_txt = q_data.get('question','')
        qtype = q_data.get('type','')
        check = q_data.get('check_point','')
        dlist = q_data.get('data',[])
        insight = st.session_state.cw_insights[i] if i < len(st.session_state.cw_insights) else ''

        slide = prs.slides.add_slide(blank_layout)

        # Q번호 뱃지
        qbadge = add_rect(slide, 2.11, 2.90, 1.93, 1.27, BRAND_RGB)
        add_textbox(slide, 2.11, 2.90, 1.93, 1.27,
                    f"Q{n}", FONT_R, 24, False, WHITE, PP_ALIGN.CENTER)

        # 질문 텍스트
        add_textbox(slide, 4.42, 2.90, 11.59, 1.27,
                    q_txt, FONT_R, 24, False, DARK)

        # 유형 뱃지
        type_badge = add_rect(slide, 2.11, 4.60, 5.82, 1.52, BRAND_RGB)
        add_textbox(slide, 2.11, 4.60, 5.82, 1.52,
                    qtype, FONT_L, 16, False, WHITE)

        # 확인 포인트
        add_textbox(slide, 8.28, 4.88, 19.94, 0.94,
                    f"확인 포인트 :  {check}", FONT_L, 16, False, GRAY)

        # 📊 라벨
        add_textbox(slide, 2.11, 7.70, 5.31, 0.94,
                    "📊  선정 근거 데이터", FONT_L, 16, False, GRAY)

        # 데이터 표
        table_data = [["출처","데이터","연도"]]
        for d in dlist:
            table_data.append([d.get('source',''), d.get('content',''), d.get('year','')])
        while len(table_data) < 4:
            table_data.append(["","",""])

        add_table_ppt(slide, 2.11, 8.82, 29.65,
                      len(table_data), 3, table_data,
                      CW_LIGHT,
                      [PPTXRGBColor(255,255,255), PPTXRGBColor(245,242,237)])

        # 인사이트 박스
        ins_text = insight if insight.strip() else f"→  이 질문에서 {st.session_state.brand_name}의 포지셔닝을 확인하세요."
        add_rect(slide, 2.11, 15.09, 29.65, 2.11, PPTXRGBColor(237,220,244))
        add_textbox(slide, 2.64, 15.72, 0.89, 0.94, "→", FONT_L, 16, False, CW_LIGHT)
        add_textbox(slide, 3.53, 15.72, 27.51, 0.94, ins_text, FONT_L, 16, False, DARK)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────
# [3] 전략제안서 Word 생성
# ─────────────────────────────────────────────
def create_proposal_word():
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.enum.table import WD_TABLE_ALIGNMENT

    FONT    = "페이퍼로지 3 Light"
    br, bg, bb = brand_rgb()
    cr, cg, cb = (112, 48, 160)
    BRAND_HEX  = st.session_state.brand_color.replace('#','')
    CW_HEX     = 'EADCF4'
    GRAY_HEX   = 'F7F7F7'
    WHITE_HEX  = 'FFFFFF'
    GREEN_HEX  = 'D9F2D0'
    ORANGE_HEX = 'FAE2D5'
    DARK_HEX   = '1A1A1A'

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Cm(1.8)
        section.bottom_margin = Cm(1.8)
        section.left_margin   = Cm(2.2)
        section.right_margin  = Cm(2.2)

    def set_bg(cell, hex_color):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_color)
        tcPr.append(shd)

    def r(para, text, size=10, bold=False, color=None, font=FONT):
        run = para.add_run(text)
        run.font.name = font
        run.font.size = Pt(size)
        run.bold = bold
        if color:
            try:
                run.font.color.rgb = RGBColor(
                    int(color[0]), int(color[1]), int(color[2]))
            except Exception:
                pass
        return run

    def add_border_para(doc, color_hex, thickness=4):
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bot = OxmlElement('w:bottom')
        bot.set(qn('w:val'), 'single')
        bot.set(qn('w:sz'), str(thickness))
        bot.set(qn('w:space'), '1')
        bot.set(qn('w:color'), color_hex)
        pBdr.append(bot)
        pPr.append(pBdr)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(4)
        return p

    def part_header(doc, part_num, title):
        p = doc.add_paragraph()
        r(p, f"PART {part_num}  ", size=16, bold=True, color=(br,bg,bb))
        r(p, title, size=14, bold=True, color=(26,23,20))
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after  = Pt(8)
        return p

    ai_keys     = ['off','on','gem']
    ai_labels   = ['GPT 검색OFF','GPT 검색ON','Gemini']

    # ── 표지 ──
    p_logo = doc.add_paragraph()
    r(p_logo, "CREAMWORKS  |  GEO 컨설팅 제안서  |  Confidential", size=9, color=(85,85,85))
    p_logo.paragraph_format.space_after = Pt(2)

    add_border_para(doc, BRAND_HEX, 6)

    p_brand = doc.add_paragraph()
    p_brand.paragraph_format.space_before = Pt(40)
    p_brand.paragraph_format.space_after  = Pt(6)
    r(p_brand, st.session_state.brand_name, size=36, color=(26,23,20))

    p_sub = doc.add_paragraph()
    r(p_sub, "AI 검색 최적화 (GEO) 전략 제안서", size=18, color=(cr,cg,cb))
    p_sub.paragraph_format.space_after = Pt(6)

    p_desc = doc.add_paragraph()
    r(p_desc, "실제 AI 답변 기반 현황 진단 + GEO 개선 전략", size=11, color=(85,85,85))
    p_desc.paragraph_format.space_after = Pt(40)

    doc.add_page_break()

    # ── 목적 박스 ──
    obj_t = doc.add_table(rows=1, cols=1)
    obj_t.style = 'Table Grid'
    set_bg(obj_t.rows[0].cells[0], CW_HEX)
    op = obj_t.rows[0].cells[0].paragraphs[0]
    r(op, "📋  이 문서의 목적  \n", size=10, bold=True, color=(cr,cg,cb))
    r(op, f"ChatGPT와 Gemini에 실제 소비자 질문 7개를 입력해 얻은 답변을 분석한 결과를 바탕으로, "
          f"{st.session_state.brand_name}의 현재 AI 검색 노출 현황을 진단하고 구체적인 GEO 개선 전략을 제시합니다.",
          size=10, color=(26,23,20))
    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # ── PART 0: 질문 설계 ──
    doc.add_page_break()
    part_header(doc, 0, "진단 질문 설계 — 근거와 방법론")

    q_table = doc.add_table(rows=len(st.session_state.questions)+1, cols=3)
    q_table.style = 'Table Grid'
    for ci, h in enumerate(["번호","단계","질문"]):
        set_bg(q_table.rows[0].cells[ci], CW_HEX)
        r(q_table.rows[0].cells[ci].paragraphs[0], h, size=9, bold=True, color=(85,85,85))
    for i, q_data in enumerate(st.session_state.questions):
        n   = i + 1
        row = q_table.rows[n]
        bg  = WHITE_HEX if i%2==0 else 'FAFAFA'
        for ci in range(3):
            set_bg(row.cells[ci], bg)
        r(row.cells[0].paragraphs[0], f"Q{n}", size=9)
        r(row.cells[1].paragraphs[0], q_data.get('stage',''), size=9)
        r(row.cells[2].paragraphs[0], q_data.get('question',''), size=9)

    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # ── PART 1: AI 진단 결과 ──
    doc.add_page_break()
    part_header(doc, 1, "AI 진단 결과 — 현황 분석")

    # 핵심 결론
    conc_t = doc.add_table(rows=1, cols=1)
    conc_t.style = 'Table Grid'
    set_bg(conc_t.rows[0].cells[0], CW_HEX)
    cp = conc_t.rows[0].cells[0].paragraphs[0]
    r(cp, "🔑  핵심 결론  \n", size=10, bold=True, color=(cr,cg,cb))
    r(cp, st.session_state.analysis_result[:400] + "...", size=10, color=(26,23,20))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # B2A 매트릭스
    mt = doc.add_table(rows=len(st.session_state.questions)+1, cols=5)
    mt.style = 'Table Grid'
    for ci, h in enumerate(["질문","내용","GPT 검색OFF","GPT 검색ON","Gemini"]):
        set_bg(mt.rows[0].cells[ci], CW_HEX)
        r(mt.rows[0].cells[ci].paragraphs[0], h, size=9, bold=True, color=(85,85,85))
    for i, q_data in enumerate(st.session_state.questions):
        n   = i + 1
        row = mt.rows[n]
        bg  = WHITE_HEX if i%2==0 else 'FEFFF0'
        for ci in range(5):
            set_bg(row.cells[ci], bg)
        r(row.cells[0].paragraphs[0], f"Q{n}", size=9)
        r(row.cells[1].paragraphs[0], q_data.get('question','')[:30], size=9)
        for j, key in enumerate(ai_keys):
            ans     = st.session_state.answers[key][n]
            mention = check_mention(ans, st.session_state.brand_name)
            cell    = row.cells[j+2]
            if ans.strip():
                if mention:
                    set_bg(cell, 'D9F2D0')
                    r(cell.paragraphs[0], "✅ 언급됨", size=9, color=(21,87,36))
                else:
                    set_bg(cell, 'FAE2D5')
                    r(cell.paragraphs[0], "❌ 미언급", size=9, color=(180,50,30))
            else:
                r(cell.paragraphs[0], "—", size=9, color=(150,150,150))

    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # ── PART 2: GEO 기초 진단 ──
    doc.add_page_break()
    part_header(doc, 2, "브랜드 GEO 기초 진단")

    gd_t = doc.add_table(rows=1, cols=2)
    gd_t.style = 'Table Grid'
    set_bg(gd_t.rows[0].cells[0], GREEN_HEX)
    set_bg(gd_t.rows[0].cells[1], ORANGE_HEX)
    r(gd_t.rows[0].cells[0].paragraphs[0],
      f"✅ GEO 강점 (AI에게 유리한 자산)\n{st.session_state.brand_usp}", size=9, color=(26,23,20))
    r(gd_t.rows[0].cells[1].paragraphs[0],
      f"❌ GEO 약점 (AI가 모르거나 약한 것)\n{st.session_state.brand_negative}", size=9, color=(26,23,20))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # 경쟁사
    competitors = [c.strip() for c in st.session_state.brand_competitors.split(',')]
    comp_t = doc.add_table(rows=len(competitors)+1, cols=3)
    comp_t.style = 'Table Grid'
    for ci, h in enumerate(["경쟁사","AI 내 현재 포지션","교촌과의 격차"]):
        set_bg(comp_t.rows[0].cells[ci], CW_HEX)
        r(comp_t.rows[0].cells[ci].paragraphs[0], h, size=9, bold=True, color=(85,85,85))
    for i, comp in enumerate(competitors):
        row = comp_t.rows[i+1]
        bg  = WHITE_HEX if i%2==0 else 'FAFAFA'
        for ci in range(3):
            set_bg(row.cells[ci], bg)
        r(row.cells[0].paragraphs[0], comp, size=9)
        r(row.cells[1].paragraphs[0], "AI 내 포지션 분석 필요", size=9, color=(150,150,150))
        r(row.cells[2].paragraphs[0], "분석 결과 참조", size=9, color=(150,150,150))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # ── PART 3~4: 전략 인사이트 ──
    doc.add_page_break()
    part_header(doc, 3, "크림웍스 GEO 전략 인사이트")

    diag_t = doc.add_table(rows=1, cols=1)
    diag_t.style = 'Table Grid'
    set_bg(diag_t.rows[0].cells[0], CW_HEX)
    dp = diag_t.rows[0].cells[0].paragraphs[0]
    r(dp, "전체 현황 진단  \n", size=10, bold=True, color=(cr,cg,cb))
    r(dp, st.session_state.overall_diagnosis or st.session_state.analysis_result[:300], size=10, color=(26,23,20))
    doc.add_paragraph().paragraph_format.space_after = Pt(4)

    for i, q_data in enumerate(st.session_state.questions):
        n       = i + 1
        insight = st.session_state.cw_insights[i]
        if insight.strip():
            ins_t = doc.add_table(rows=2, cols=1)
            ins_t.style = 'Table Grid'
            set_bg(ins_t.rows[0].cells[0], DARK_HEX)
            r(ins_t.rows[0].cells[0].paragraphs[0],
              f"Q{n}. {q_data.get('question','')}", size=10, bold=True, color=(br,bg,bb))
            set_bg(ins_t.rows[1].cells[0], CW_HEX)
            r(ins_t.rows[1].cells[0].paragraphs[0],
              f"💜 크림웍스 전략: {insight}", size=10, color=(cr,cg,cb))
            doc.add_paragraph().paragraph_format.space_after = Pt(4)

    # ── PART 4: 액션 플랜 ──
    doc.add_page_break()
    part_header(doc, 4, "우선 실행 액션 플랜")

    action_t = doc.add_table(rows=3, cols=2)
    action_t.style = 'Table Grid'
    for ri, (label, content) in enumerate([
        ("🚨 즉시 실행", ""),
        ("⭐ 1개월 내", ""),
        ("3개월 내", ""),
    ]):
        set_bg(action_t.rows[ri].cells[0], CW_HEX if ri>0 else DARK_HEX)
        color = (br,bg,bb) if ri==0 else (cr,cg,cb)
        r(action_t.rows[ri].cells[0].paragraphs[0], label, size=10, bold=True, color=color)
        r(action_t.rows[ri].cells[1].paragraphs[0], content or "—", size=10)

    lines = st.session_state.priority_actions.split('\n')
    for line in lines:
        if line.strip():
            p = doc.add_paragraph()
            r(p, line.strip(), size=10, color=(26,23,20))

    # 푸터
    fp = doc.add_paragraph()
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r(fp, "CREAMWORKS  —  AI가 좋아하는 브랜드를 만듭니다", size=9, color=(150,150,150))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────
# [4] 전략제안서 PPT 생성
# ─────────────────────────────────────────────
def create_proposal_ppt():
    FONT_M  = "페이퍼로지 5 Medium"
    FONT_R  = "페이퍼로지 4 Regular"
    FONT_L  = "페이퍼로지 3 Light"
    br, bg, bb = brand_rgb()
    BRAND_RGB = PPTXRGBColor(br, bg, bb)
    CW_PURPLE = PPTXRGBColor(83, 39, 168)
    CW_LIGHT  = PPTXRGBColor(124, 92, 191)
    DARK      = PPTXRGBColor(26, 23, 20)
    GRAY      = PPTXRGBColor(85, 85, 85)
    WHITE     = PPTXRGBColor(255, 255, 255)
    GOLD      = BRAND_RGB   # PART 번호에 브랜드 컬러 적용
    GREEN     = PPTXRGBColor(217, 242, 208)
    ORANGE_C  = PPTXRGBColor(250, 226, 213)
    CW_LAVENDER = PPTXRGBColor(234, 220, 244)

    W = Inches(13.33)
    H = Inches(7.50)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    ai_keys   = ['off','on','gem']
    ai_labels = ['GPT 검색OFF','GPT 검색ON','Gemini']

    def add_textbox(slide, left, top, width, height, text, font_name=FONT_L,
                    font_size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(
            PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = PPTXPt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return tb

    def add_rect(slide, left, top, width, height, fill_color):
        shape = slide.shapes.add_shape(
            1, PPTXCm(left), PPTXCm(top), PPTXCm(width), PPTXCm(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def part_slide(title_num, title_text):
        slide = prs.slides.add_slide(blank)
        add_textbox(slide, 2.82, 6.68, 6.05, 2.13,
                    f"PART {title_num}", FONT_R, 44, False, GOLD)
        add_rect(slide, 2.82, 8.81, 4.24, 0.15, GOLD)
        add_textbox(slide, 2.82, 9.22, 12.0, 1.27,
                    title_text, FONT_R, 24, False, DARK)
        return slide

    # ── 슬라이드 1: 표지 ──
    slide1 = prs.slides.add_slide(blank)
    add_rect(slide1, 0, 0, 33.87, 19.05, PPTXRGBColor(245,245,245))
    add_rect(slide1, 8.56, 4.55, 0.46, 5.39, CW_PURPLE)
    add_textbox(slide1, 9.38, 4.73, 12.0, 1.81,
                "크림웍스\nGEO 컨설팅", FONT_M, 36, False, DARK)
    add_textbox(slide1, 9.38, 6.55, 7.0, 1.27,
                "[3. 전략 제안서]", FONT_R, 24, False, DARK)
    add_textbox(slide1, 9.38, 7.73, 14.1, 1.63,
                "AI 검색시대, 브랜드가 ChatGPT·Gemini 대답에서\n발견된 브랜드의 현재 상태에 대한 보고서입니다.",
                FONT_R, 16, False, DARK)
    add_textbox(slide1, 11.20, 17.69, 11.47, 0.86,
                "CREAMWORKS  - AI가 좋아하는 브랜드를 만듭니다",
                FONT_M, 14, False, CW_PURPLE)

    # ── PART 0 섹션 ──
    part_slide(0, "진단 질문 설계 - 근거와 방법론")

    # ── PART 0 내용 슬라이드 ──
    s0 = prs.slides.add_slide(blank)
    add_textbox(s0, 2.11, 1.98, 11.51, 1.45,
                "0. 설계 철학 및 질문별 개요", FONT_R, 28, False, DARK)
    add_textbox(s0, 2.13, 3.89, 29.65, 1.80,
                f"브랜드를 알던 모르던 일반적 소비자가 AI에게 정보를 요청하는 순간을 포착하는 것이 GEO의 출발점입니다.\n"
                f"이 철학을 기반으로 DISCOVER → CONSIDER → DECIDE 여정 순서로 질문 7개가 설계되었습니다.",
                FONT_L, 14, False, DARK)

    # 질문 목록 표
    q_rows = [["번호","단계","질문 내용"]]
    for i, q_data in enumerate(st.session_state.questions):
        q_rows.append([f"Q{i+1}", q_data.get('stage',''), q_data.get('question','')])
    tbl = s0.shapes.add_table(len(q_rows), 3,
        PPTXCm(2.11), PPTXCm(6.12), PPTXCm(29.65), PPTXCm(11.10)).table
    tbl.columns[0].width = PPTXCm(2.5)
    tbl.columns[1].width = PPTXCm(5.0)
    tbl.columns[2].width = PPTXCm(22.15)
    for ri, row_data in enumerate(q_rows):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = CW_LIGHT
                clr = WHITE
                bold = True
            else:
                cell.fill.fore_color.rgb = PPTXRGBColor(255,255,255) if ri%2==1 else PPTXRGBColor(245,242,237)
                clr = DARK
                bold = False
            tf = cell.text_frame
            p  = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name  = FONT_L
            run.font.size  = PPTXPt(13)
            run.font.bold  = bold
            run.font.color.rgb = clr

    # ── PART 1 섹션 ──
    part_slide(1, "AI 진단결과 - 현황분석")

    # PART 1 내용: B2A 매트릭스
    s1 = prs.slides.add_slide(blank)
    add_textbox(s1, 2.11, 0.56, 15.0, 1.12,
                "1. B2A 매트릭스 — AI 답변 기반 언급 현황", FONT_R, 22, False, DARK)

    b2a_rows = [["질문","GPT 검색OFF","GPT 검색ON","Gemini"]]
    for i, q_data in enumerate(st.session_state.questions):
        n = i + 1
        row_data = [f"Q{n}. {q_data.get('question','')[:25]}"]
        for key in ai_keys:
            ans     = st.session_state.answers[key][n]
            mention = check_mention(ans, st.session_state.brand_name)
            if ans.strip():
                row_data.append("✅ 언급됨" if mention else "❌ 미언급")
            else:
                row_data.append("—")
        b2a_rows.append(row_data)

    tbl2 = s1.shapes.add_table(len(b2a_rows), 4,
        PPTXCm(2.11), PPTXCm(2.0), PPTXCm(29.65), PPTXCm(16.0)).table
    tbl2.columns[0].width = PPTXCm(12.0)
    for ci in range(1, 4):
        tbl2.columns[ci].width = PPTXCm(5.88)
    for ri, row_data in enumerate(b2a_rows):
        for ci, val in enumerate(row_data):
            cell = tbl2.cell(ri, ci)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = CW_LIGHT
                clr = WHITE
                bold = True
            else:
                if "✅" in val:
                    cell.fill.fore_color.rgb = PPTXRGBColor(217,242,208)
                    clr = PPTXRGBColor(21,87,36)
                elif "❌" in val:
                    cell.fill.fore_color.rgb = PPTXRGBColor(250,226,213)
                    clr = PPTXRGBColor(180,50,30)
                else:
                    cell.fill.fore_color.rgb = PPTXRGBColor(255,255,255) if ri%2==1 else PPTXRGBColor(245,242,237)
                    clr = GRAY
                bold = False
            tf = cell.text_frame
            p  = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name  = FONT_L
            run.font.size  = PPTXPt(13)
            run.font.bold  = bold
            run.font.color.rgb = clr

    # ── PART 2 섹션 ──
    part_slide(2, "브랜드 GEO 기초 진단")

    # ── PART 3 섹션 ──
    part_slide(3, "크림웍스 GEO 전략 인사이트")

    # PART 3 내용: 인사이트 슬라이드들
    for i, q_data in enumerate(st.session_state.questions):
        n       = i + 1
        insight = st.session_state.cw_insights[i]
        if not insight.strip():
            continue
        si = prs.slides.add_slide(blank)
        add_rect(si, 2.11, 0.56, 29.65, 1.52, CW_LIGHT)
        add_textbox(si, 2.64, 0.79, 28.0, 1.02,
                    f"Q{n}. {q_data.get('question','')}", FONT_R, 20, True, WHITE)
        add_rect(si, 2.11, 2.41, 29.65, 5.46, CW_LAVENDER)
        add_textbox(si, 2.64, 2.79, 0.89, 0.94, "→", FONT_L, 20, False, CW_LIGHT)
        add_textbox(si, 3.81, 2.79, 27.51, 4.37, insight, FONT_L, 16, False, DARK)

    # ── PART 4 섹션 ──
    part_slide(4, "우선 실행 액션 플랜")

    # 액션 슬라이드
    sa = prs.slides.add_slide(blank)
    add_textbox(sa, 2.11, 0.56, 15.0, 1.12,
                "4. GEO 실행 액션 플랜", FONT_R, 22, False, DARK)

    actions = [
        ("🚨 즉시", "robots.txt, Schema Markup, 브랜드명 통일 표기"),
        ("⭐ 1개월", "FAQ 콘텐츠 7개, llms.txt, 나무위키 보강"),
        ("3개월", "언론 PR, 채널별 콘텐츠 배포, B2A 월간 측정 시작"),
    ]
    for ai, (label, content) in enumerate(actions):
        top = 2.2 + ai * 4.5
        add_rect(sa, 2.11, top, 29.65, 3.81, CW_LAVENDER if ai>0 else PPTXRGBColor(26,23,20))
        clr = BRAND_RGB if ai==0 else CW_PURPLE
        add_textbox(sa, 2.64, top+0.51, 5.0, 0.94, label, FONT_R, 20, True, clr)
        add_textbox(sa, 8.0, top+0.51, 23.0, 2.54, content, FONT_L, 16, False, DARK)

    # 마지막: 푸터 슬라이드
    sf = prs.slides.add_slide(blank)
    add_rect(sf, 0, 0, 33.87, 19.05, PPTXRGBColor(26,23,20))
    add_textbox(sf, 9.38, 8.38, 15.0, 1.45,
                "CREAMWORKS  —  AI가 좋아하는 브랜드를 만듭니다",
                FONT_M, 20, False, CW_PURPLE, PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf




import json
from datetime import datetime

# ─────────────────────────────────────────────
# 페이지 설정
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="GEO-Scan | 크림웍스",
    layout="wide",
    initial_sidebar_state="collapsed"
)

CW_PURPLE = "#6B4EFF"
CW_PURPLE_LIGHT = "#EDE9FF"
CW_PURPLE_RGB = (107, 78, 255)

st.markdown(f"""
<style>
  .main .block-container {{ padding-top: 1.5rem; max-width: 860px; margin: 0 auto; }}
  .brand-header {{
      padding: 24px 28px;
      border-radius: 14px;
      margin-bottom: 1.5rem;
      color: white;
  }}
  .step-label {{
      font-size: 0.72rem;
      font-weight: 600;
      color: #999;
      letter-spacing: 1px;
      text-transform: uppercase;
      margin-bottom: 2px;
  }}
  .step-title {{
      font-size: 1.1rem;
      font-weight: 700;
      color: #1a1a1a;
      margin-bottom: 0.8rem;
  }}
  .field-label {{
      font-size: 0.82rem;
      font-weight: 600;
      color: #444;
      margin-bottom: 4px;
      margin-top: 10px;
  }}
  .required {{ color: #e74c3c; }}
  .cw-box {{
      background: {CW_PURPLE_LIGHT};
      border-left: 4px solid {CW_PURPLE};
      padding: 14px 18px;
      border-radius: 0 10px 10px 0;
      margin: 8px 0 16px 0;
      font-size: 0.88rem;
      color: #3a2d7a;
  }}
  .info-box {{
      background: #f0f4ff;
      border: 1px solid #c7d3ff;
      border-radius: 10px;
      padding: 14px 18px;
      font-size: 0.88rem;
      color: #2c3e80;
      margin-bottom: 1rem;
  }}
  .mention-yes {{
      background: #d4edda; color: #155724;
      padding: 3px 10px; border-radius: 10px;
      font-size: 0.78rem; font-weight: 700;
  }}
  .mention-no {{
      background: #f8d7da; color: #721c24;
      padding: 3px 10px; border-radius: 10px;
      font-size: 0.78rem; font-weight: 700;
  }}
  .q-card {{
      border: 1px solid #e8e8e8;
      border-radius: 10px;
      padding: 16px 18px;
      margin-bottom: 12px;
      background: white;
  }}
  .q-num {{
      display: inline-block;
      background: #111;
      color: white;
      width: 22px; height: 22px;
      border-radius: 50%;
      text-align: center;
      line-height: 22px;
      font-size: 0.75rem;
      font-weight: 700;
      margin-right: 8px;
  }}
  .divider {{ border: none; border-top: 1px solid #eee; margin: 1.5rem 0; }}

  /* 연두색 액션 버튼 */
  div[data-testid="stButton"] button[kind="primary"] {{
      background-color: #52B788 !important;
      border-color: #52B788 !important;
      color: white !important;
      font-weight: 600 !important;
      max-width: 320px !important;
      border-radius: 8px !important;
  }}
  div[data-testid="stButton"] button[kind="primary"]:hover {{
      background-color: #40916C !important;
      border-color: #40916C !important;
  }}
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────
# 세션 초기화
# ─────────────────────────────────────────────
def init():
    defaults = {
        'step': 1,
        'api_key': '',
        'brand_name': '',
        'brand_color': '#4A90D9',
        'brand_category': '',
        'brand_usp': '',
        'brand_target': '',
        'brand_competitors': '',
        'brand_negative': '',
        'brand_focus': '',
        'questions': [],
        'questions_confirmed': False,
        'answers': {
            'off': {i: '' for i in range(1, 8)},
            'on':  {i: '' for i in range(1, 8)},
            'gem': {i: '' for i in range(1, 8)},
        },
        'analysis_result': '',
        'cw_insights': [''] * 7,
        'overall_diagnosis': '',
        'priority_actions': '🚨 즉시 실행:\n\n⭐ 1개월 내:\n\n3개월 내:',
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

init()

# ─────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────
def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color.replace('#', ''))
    tcPr.append(shd)

def check_mention(text, brand_name):
    if not text.strip():
        return None
    keywords = [brand_name, brand_name[:2], brand_name.replace(' ', '')]
    return any(kw in text for kw in keywords)

def get_client():
    return anthropic.Anthropic(api_key=st.session_state.api_key)

# ─────────────────────────────────────────────
# 헤더
# ─────────────────────────────────────────────
bc = st.session_state.brand_color
bn = st.session_state.brand_name or 'GEO-Scan'
st.markdown(f"""
<div class="brand-header" style="background: linear-gradient(135deg, #0f0f0f 60%, {bc}99);">
  <div style="font-size:0.78rem;color:#888;letter-spacing:2px;font-weight:600;">CREAMWORKS  ·  GEO-Scan v2.0</div>
  <div style="font-size:1.7rem;font-weight:800;margin:6px 0 4px;">{bn} AI 진단 시스템</div>
  <div style="font-size:0.85rem;color:#bbb;">브랜드 AI 검색 노출 현황 진단 · GEO 전략 제안서 자동 생성</div>
</div>
""", unsafe_allow_html=True)

# 스텝 인디케이터
step_names = ["브랜드 입력", "질문 확인", "답변 수집", "분석·인사이트", "보고서"]
cols = st.columns(5)
for i, (col, name) in enumerate(zip(cols, step_names)):
    n = i + 1
    active = st.session_state.step == n
    done = st.session_state.step > n
    bg = CW_PURPLE if active else ('#27ae60' if done else '#ddd')
    tc = 'white'
    label_c = '#1a1a1a' if active else ('#27ae60' if done else '#aaa')
    col.markdown(f"""
    <div style="text-align:center;padding:4px 0">
      <div style="width:26px;height:26px;border-radius:50%;background:{bg};
                  color:{tc};display:inline-flex;align-items:center;justify-content:center;
                  font-size:0.72rem;font-weight:700;margin-bottom:3px">
        {'✓' if done else n}
      </div>
      <div style="font-size:0.7rem;color:{label_c};font-weight:{'700' if active else '400'}">{name}</div>
    </div>""", unsafe_allow_html=True)

st.markdown("<hr class='divider'>", unsafe_allow_html=True)

# ─────────────────────────────────────────────
# 이전 작업 이어하기 (Excel 업로드)
# ─────────────────────────────────────────────
if st.session_state.step == 1:
    with st.expander("📂 이전 작업 이어하기 (Excel 업로드)", expanded=False):
        st.caption("이전에 저장한 Excel 파일을 업로드하면 중간부터 이어서 진행할 수 있습니다.")

        col_a, col_b = st.columns(2)
        with col_a:
            st.markdown("**케이스 A — 답변수집 Excel**")
            st.caption("브랜드 정보 + 질문 + 답변 복원 → STEP 4 (분석) 부터 시작")
        with col_b:
            st.markdown("**케이스 B — 전체데이터 Excel**")
            st.caption("브랜드 정보 + 질문 + 답변 + 분석 복원 → STEP 5 (보고서) 부터 시작")

        uploaded_xl = st.file_uploader(
            "Excel 파일 업로드 (답변수집 또는 전체데이터)",
            type=["xlsx"],
            key="resume_upload"
        )

        if uploaded_xl is not None:
            try:
                file_type = load_from_excel(uploaded_xl)
                if file_type == 'answer':
                    st.success(f"✅ **답변수집 파일** 복원 완료! → STEP 4 (분석 시작) 로 이동합니다.")
                    brand = st.session_state.get('brand_name','')
                    st.info(f"브랜드: **{brand}** | 질문 **{len(st.session_state.questions)}개** | 답변 복원 완료")
                else:
                    st.success(f"✅ **전체데이터 파일** 복원 완료! → STEP 5 (보고서 생성) 로 이동합니다.")
                    brand = st.session_state.get('brand_name','')
                    st.info(f"브랜드: **{brand}** | Claude 분석 결과 포함")

                if st.button("▶ 이어서 진행하기", type="primary"):
                    st.rerun()

            except Exception as e:
                st.error(f"파일 읽기 오류: {e}\n올바른 GEO-Scan Excel 파일을 업로드해주세요.")

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

# ─────────────────────────────────────────────
# STEP 1: 브랜드 정보 입력
# ─────────────────────────────────────────────
if st.session_state.step == 1:
    st.markdown('<div class="step-label">STEP 1</div>', unsafe_allow_html=True)
    st.markdown('<div class="step-title">브랜드 정보 입력</div>', unsafe_allow_html=True)

    st.markdown('<div class="cw-box">💜 많이 입력할수록 질문 퀄리티가 좋아집니다. 핵심 USP와 강조 포인트는 꼭 입력해주세요.</div>', unsafe_allow_html=True)

    # API 키
    st.markdown('<div class="field-label">🔑 Anthropic API Key <span class="required">*</span></div>', unsafe_allow_html=True)
    api_key = st.text_input("api", value=st.session_state.api_key,
                             type="password", placeholder="sk-ant-api03-...",
                             label_visibility="collapsed")
    st.session_state.api_key = api_key
    st.caption("API 키는 브라우저 세션에만 유지되며 저장되지 않습니다.")

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    col1, col2 = st.columns(2)

    with col1:
        st.markdown('<div class="field-label">브랜드명 <span class="required">*</span></div>', unsafe_allow_html=True)
        st.session_state.brand_name = st.text_input(
            "브랜드명", value=st.session_state.brand_name,
            placeholder="예: 교촌치킨", label_visibility="collapsed")

        st.markdown('<div class="field-label">카테고리 <span class="required">*</span></div>', unsafe_allow_html=True)
        st.session_state.brand_category = st.text_input(
            "카테고리", value=st.session_state.brand_category,
            placeholder="예: 치킨 프랜차이즈", label_visibility="collapsed")

        st.markdown('<div class="field-label">브랜드 공식 컬러</div>', unsafe_allow_html=True)
        st.session_state.brand_color = st.color_picker(
            "컬러", value=st.session_state.brand_color, label_visibility="collapsed")
        bc_preview = st.session_state.brand_color
        st.markdown(f"""
        <div style="display:flex;align-items:center;gap:10px;margin-top:6px;margin-bottom:2px">
          <div style="width:32px;height:32px;background:{bc_preview};
                      border-radius:6px;border:1px solid #ddd;"></div>
          <span style="font-size:0.82rem;color:#666">{bc_preview} — 보고서·질문지에 자동 적용됩니다</span>
        </div>
        <div style="font-size:0.78rem;color:#aaa;margin-bottom:4px">
          구글에 "{st.session_state.brand_name or '브랜드명'} 브랜드 컬러 HEX" 검색 후 입력하세요
        </div>
        <div style="font-size:0.78rem;color:#aaa;margin-bottom:8px">
          예) 교촌치킨 → #F9BA15 &nbsp;|&nbsp; 드시모네 → #01C49C &nbsp;|&nbsp; 크림웍스 → #7C5CBF
        </div>
        """, unsafe_allow_html=True)

        st.markdown('<div class="field-label">경쟁 브랜드</div>', unsafe_allow_html=True)
        st.session_state.brand_competitors = st.text_input(
            "경쟁사", value=st.session_state.brand_competitors,
            placeholder="예: BBQ, BHC", label_visibility="collapsed")

        st.markdown('<div class="field-label">주요 타겟 소비자</div>', unsafe_allow_html=True)
        st.session_state.brand_target = st.text_input(
            "타겟", value=st.session_state.brand_target,
            placeholder="예: 남녀노소, 배달앱 구매 중심", label_visibility="collapsed")

    with col2:
        st.markdown('<div class="field-label">핵심 USP (차별점) <span class="required">*</span></div>', unsafe_allow_html=True)
        st.session_state.brand_usp = st.text_area(
            "USP", value=st.session_state.brand_usp, height=110,
            placeholder="예: 간장치킨 원조, 35년 업력, 비밀 마늘간장 소스, 붓질 공정, 국내산 재료",
            label_visibility="collapsed")

        st.markdown('<div class="field-label">부정 이미지 / 약점</div>', unsafe_allow_html=True)
        st.session_state.brand_negative = st.text_area(
            "약점", value=st.session_state.brand_negative, height=80,
            placeholder="예: 가격 인상 선도, 배달 이중가격 논란",
            label_visibility="collapsed")

        st.markdown('<div class="field-label">지금 강조하고 싶은 포인트 <span class="required">*</span></div>', unsafe_allow_html=True)
        st.session_state.brand_focus = st.text_area(
            "포인트", value=st.session_state.brand_focus, height=80,
            placeholder="예: 소비자가 교촌을 선택하는 이유 부각, 앱 구매 중심 GEO 세팅 방향",
            label_visibility="collapsed")

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    col_btn, col_empty = st.columns([2, 3])
    with col_btn:
        clicked = st.button("🔍  질문 7개 생성하기", type="primary", use_container_width=True)

    if clicked:
        # 유효성 검사
        if not st.session_state.api_key:
            st.error("Anthropic API Key를 입력해주세요.")
        elif not st.session_state.brand_name:
            st.error("브랜드명을 입력해주세요.")
        elif not st.session_state.brand_category:
            st.error("카테고리를 입력해주세요.")
        elif not st.session_state.brand_usp:
            st.error("핵심 USP를 입력해주세요.")
        else:
            with st.spinner("Claude가 브랜드를 분석하고 질문을 설계하는 중입니다..."):
                try:
                    client = get_client()

                    prompt = f"""당신은 GEO(Generative Engine Optimization) 전문가입니다.
아래 브랜드 정보를 바탕으로, 소비자가 ChatGPT·Gemini에 실제로 물어볼 법한 GEO 진단 질문 7개를 설계해주세요.

브랜드 정보:
- 브랜드명: {st.session_state.brand_name}
- 카테고리: {st.session_state.brand_category}
- 핵심 USP: {st.session_state.brand_usp}
- 주요 타겟: {st.session_state.brand_target}
- 경쟁 브랜드: {st.session_state.brand_competitors}
- 부정 이미지/약점: {st.session_state.brand_negative}
- 강조 포인트: {st.session_state.brand_focus}

질문 설계 원칙:
1. 브랜드명이 절대 들어가면 안 됨 (브랜드를 모르는 소비자가 AI에게 묻는 질문)
2. 실제 소비자가 쓰는 구어체
3. AIJ 단계: DISCOVER(2개) / CONSIDER(3개) / DECIDE(2개)
4. 브랜드 USP와 직접 연결되는 질문 우선
5. 부정 이미지 방어 질문 1개 이상 포함

⚠️ data 필드 작성 필수 규칙:
- 각 질문마다 반드시 실제 존재하는 기관/매체의 데이터 3개를 넣어야 함
- "출처기관" 같은 템플릿 문자열 절대 금지. 반드시 실제 기관명 사용
- 실제 기관 예시: 한국소비자원, 식품의약품안전처, 오픈서베이, 네이버 데이터랩, 닐슨코리아, 건강보험심사평가원, 통계청, aT한국농수산식품유통공사, 매일경제, 조선일보, 와이즈앱, 서울대병원, 대한의사협회 등 (카테고리에 맞는 기관 선택)
- content는 해당 카테고리와 직접 관련된 구체적 수치/통계/트렌드 포함
- year는 2023~2026 사이 실제 연도
- check_point는 구체적으로: "{st.session_state.brand_name}이(가) [구체적 맥락]으로 등장하는지 + [{st.session_state.brand_competitors}] 대비 [포지션] 언급 여부" 형식

반드시 아래 JSON 형식으로만 응답 (다른 텍스트 없이):
{{
  "questions": [
    {{
      "question": "질문 내용 (구어체, 브랜드명 없이)",
      "stage": "DISCOVER 또는 CONSIDER 또는 DECIDE",
      "type": "유형명 (예: 카테고리 진입 — 브랜드 선택 첫 질문)",
      "check_point": "확인 포인트: {st.session_state.brand_name}이(가) [구체적 맥락]으로 등장하는지 + [{st.session_state.brand_competitors}] 대비 포지션 언급 여부",
      "data": [
        {{"source": "실제기관명", "content": "구체적 수치 포함한 데이터 내용", "year": "2024"}},
        {{"source": "실제기관명", "content": "구체적 수치 포함한 데이터 내용", "year": "2025"}},
        {{"source": "실제기관명", "content": "구체적 수치 포함한 데이터 내용", "year": "2024"}}
      ]
    }}
  ]
}}"""

                    message = client.messages.create(
                        model="claude-sonnet-4-20250514",
                        max_tokens=8000,
                        messages=[{"role": "user", "content": prompt}]
                    )

                    raw = message.content[0].text.strip()
                    if "```" in raw:
                        raw = raw.split("```")[1]
                        if raw.startswith("json"):
                            raw = raw[4:]
                    raw = raw.strip()

                    data = json.loads(raw)
                    st.session_state.questions = data['questions']
                    st.session_state.questions_confirmed = False
                    st.session_state.step = 2
                    st.rerun()

                except json.JSONDecodeError:
                    st.error("질문 생성 중 파싱 오류가 발생했습니다. 다시 시도해주세요.")
                except Exception as e:
                    st.error(f"오류 발생: {str(e)}")

# ─────────────────────────────────────────────
# STEP 2: 질문 확인 및 수정
# ─────────────────────────────────────────────
elif st.session_state.step == 2:
    st.markdown('<div class="step-label">STEP 2</div>', unsafe_allow_html=True)
    st.markdown('<div class="step-title">진단 질문 7개 확인 및 수정</div>', unsafe_allow_html=True)
    st.markdown('<div class="cw-box">💜 Claude가 설계한 질문입니다. 수정이 필요하면 직접 편집하세요. 브랜드명이 포함된 질문은 사용하지 않습니다.</div>', unsafe_allow_html=True)

    updated_qs = []
    for i, q_data in enumerate(st.session_state.questions):
        n = i + 1
        stage = q_data.get('stage', '')
        qtype = q_data.get('type', '')
        check = q_data.get('check_point', '')
        data_list = q_data.get('data', [])

        with st.expander(f"Q{n} · [{stage}] {qtype}", expanded=(i < 2)):
            new_q = st.text_input(
                f"질문 내용",
                value=q_data.get('question', ''),
                key=f"q_edit_{i}",
                label_visibility="visible"
            )
            st.caption(f"**확인 포인트:** {check}")

            if data_list:
                st.markdown("**📊 선정 근거 데이터**")
                cols = st.columns(3)
                for j, d in enumerate(data_list[:3]):
                    with cols[j]:
                        st.markdown(f"""
                        <div style="background:#f8f9fa;border-radius:8px;padding:10px;font-size:0.78rem;">
                          <div style="font-weight:700;color:#555;margin-bottom:4px">{d.get('source','')} ({d.get('year','')})</div>
                          <div style="color:#333">{d.get('content','')}</div>
                        </div>""", unsafe_allow_html=True)

            q_updated = dict(q_data)
            q_updated['question'] = new_q
            updated_qs.append(q_updated)

    col_back, col_confirm, col_empty2 = st.columns([1, 2, 2])
    with col_back:
        if st.button("← 다시 입력", use_container_width=True):
            st.session_state.step = 1
            st.rerun()
    with col_confirm:
        if st.button("✅ 질문 확정 — 답변 수집으로 이동", type="primary", use_container_width=True):
            st.session_state.questions = updated_qs
            st.session_state.questions_confirmed = True
            st.session_state.answers = {
                'off': {i: '' for i in range(1, 8)},
                'on':  {i: '' for i in range(1, 8)},
                'gem': {i: '' for i in range(1, 8)},
            }
            st.session_state.step = 3
            st.rerun()

# ─────────────────────────────────────────────
# STEP 3: 답변 수집
# ─────────────────────────────────────────────
elif st.session_state.step == 3:
    st.markdown('<div class="step-label">STEP 3</div>', unsafe_allow_html=True)
    st.markdown('<div class="step-title">AI 답변 수집</div>', unsafe_allow_html=True)

    # 진행 현황
    off_cnt = sum(1 for v in st.session_state.answers['off'].values() if v.strip())
    on_cnt  = sum(1 for v in st.session_state.answers['on'].values() if v.strip())
    gem_cnt = sum(1 for v in st.session_state.answers['gem'].values() if v.strip())
    total   = off_cnt + on_cnt + gem_cnt

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("GPT 검색OFF", f"{off_cnt}/7")
    c2.metric("GPT 검색ON", f"{on_cnt}/7")
    c3.metric("Gemini", f"{gem_cnt}/7")
    c4.metric("전체", f"{total}/21")
    st.progress(total / 21)

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    ai_tabs = st.tabs(["🔵 ChatGPT 검색OFF", "🟢 ChatGPT 검색ON", "🟠 Gemini"])
    ai_keys   = ['off', 'on', 'gem']
    ai_labels = ['ChatGPT 검색OFF', 'ChatGPT 검색ON', 'Gemini']
    ai_instrs = [
        "ChatGPT → 새 채팅 → 메모리 OFF + 검색 OFF → Q1~Q7 순서대로 동일 채팅 입력 → 각 답변 복사 후 붙여넣기",
        "ChatGPT → 새 채팅 → 메모리 OFF + 검색 ON → Q1~Q7 순서대로 동일 채팅 입력 → 각 답변 복사 후 붙여넣기",
        "시크릿 모드 → gemini.google.com (로그아웃) → 새 채팅 → Q1~Q7 순서대로 입력 → 각 답변 복사 후 붙여넣기"
    ]

    for ai_tab, key, label, instr in zip(ai_tabs, ai_keys, ai_labels, ai_instrs):
        with ai_tab:
            st.markdown(f'<div class="info-box">📌 {instr}</div>', unsafe_allow_html=True)

            for i, q_data in enumerate(st.session_state.questions):
                n = i + 1
                q_text = q_data.get('question', '')
                ans_key = f"ans_{key}_{n}"

                # on_change 콜백으로 즉시 반영
                def make_callback(k, num):
                    def _cb():
                        val = st.session_state.get(f"ans_{k}_{num}", "")
                        st.session_state.answers[k][num] = val
                    return _cb

                # 현재 저장 상태 확인
                current_val = st.session_state.get(ans_key, st.session_state.answers[key][n])
                saved = bool(current_val.strip())

                col_q, col_badge = st.columns([5, 1])
                with col_q:
                    st.markdown(f'<div style="margin:8px 0 2px;font-weight:600">Q{n}.</div>', unsafe_allow_html=True)
                    st.code(q_text, language=None)
                with col_badge:
                    st.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)
                    if saved:
                        st.markdown('<span class="mention-yes">✓ 저장됨</span>', unsafe_allow_html=True)
                    else:
                        st.markdown('<span class="mention-no">미입력</span>', unsafe_allow_html=True)

                st.text_area(
                    f"답변",
                    value=st.session_state.answers[key][n],
                    key=ans_key,
                    height=110,
                    placeholder=f"{label} 답변을 여기에 붙여넣으세요...",
                    label_visibility="collapsed",
                    on_change=make_callback(key, n)
                )
                st.markdown("<hr style='border:none;border-top:1px solid #f0f0f0;margin:4px 0'>", unsafe_allow_html=True)

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    # Excel 저장 버튼
    st.markdown("#### 💾 답변 저장")
    st.caption("수집한 답변을 Excel 파일로 저장해두세요. 세션이 종료되면 데이터가 사라집니다.")

    if st.button("📊 답변 Excel 저장", use_container_width=True):
        if total == 0:
            st.error("저장할 답변이 없습니다. 먼저 답변을 입력해주세요.")
        else:
            with st.spinner("Excel 파일 생성 중..."):
                buf = create_answer_excel()
            fname = f"{st.session_state.brand_name}_GEO_답변수집_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
            st.download_button(
                label=f"⬇️ {fname} 다운로드",
                data=buf,
                file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
                key="dl_excel_ans"
            )
            st.success(f"Excel 파일이 생성되었습니다! ({total}/21개 답변 저장)")

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)
    col_back, col_next, col_empty3 = st.columns([1, 2, 2])
    with col_back:
        if st.button("← 질문 수정", use_container_width=True):
            st.session_state.step = 2
            st.rerun()
    with col_next:
        if st.button("📊 분석 시작 →", type="primary", use_container_width=True):
            if total == 0:
                st.error("최소 1개 이상의 답변을 입력해주세요.")
            else:
                # Claude API로 B2A 분석
                with st.spinner("Claude가 21개 답변을 분석 중입니다..."):
                    try:
                        client = get_client()
                        all_answers = []
                        for i, q_data in enumerate(st.session_state.questions):
                            n = i + 1
                            q_text = q_data.get('question', '')
                            for key, label in zip(ai_keys, ai_labels):
                                ans = st.session_state.answers[key][n]
                                if ans.strip():
                                    all_answers.append(f"[{label}] Q{n}: {q_text}\n답변: {ans}")

                        analysis_prompt = f"""당신은 GEO 분석 전문가입니다.
브랜드: {st.session_state.brand_name}
경쟁사: {st.session_state.brand_competitors}

아래는 ChatGPT(검색OFF/ON)와 Gemini에 실제 소비자 질문을 입력해서 얻은 AI 답변입니다.
각 답변을 분석해서 아래 형식으로 정리해주세요.

분석 항목:
1. 전체 요약: {st.session_state.brand_name}의 현재 AI 노출 현황 핵심 요약 (3~5문장)
2. 질문별 분석: Q1~Q7 각각에 대해
   - 브랜드 언급 여부 (O/X)
   - 언급 맥락 (긍정/중립/부정, 몇 번째 언급)
   - 경쟁사 언급 현황
   - 핵심 발견사항
3. 경쟁사 포지션: AI가 경쟁사를 어떻게 포지셔닝하는지
4. 핵심 공백: {st.session_state.brand_name}이 AI에서 완전히 빠진 영역
5. GEO 기회: 가장 빠르게 개선 가능한 포인트 3가지

--- 수집된 답변 ---
{chr(10).join(all_answers)}
"""
                        message = client.messages.create(
                            model="claude-sonnet-4-20250514",
                            max_tokens=3000,
                            messages=[{"role": "user", "content": analysis_prompt}]
                        )
                        st.session_state.analysis_result = message.content[0].text
                        st.session_state.cw_insights = [''] * 7
                        st.session_state.step = 4
                        st.rerun()

                    except Exception as e:
                        st.error(f"분석 오류: {str(e)}")

# ─────────────────────────────────────────────
# STEP 4: 분석 결과 + 인사이트
# ─────────────────────────────────────────────
elif st.session_state.step == 4:
    st.markdown('<div class="step-label">STEP 4</div>', unsafe_allow_html=True)
    st.markdown('<div class="step-title">B2A 분석 결과 + 크림웍스 전략 인사이트</div>', unsafe_allow_html=True)

    # B2A 매트릭스
    st.markdown("#### 📊 B2A 매트릭스 — 언급 현황")
    ai_keys   = ['off', 'on', 'gem']
    ai_labels_short = ['GPT 검색OFF', 'GPT 검색ON', 'Gemini']

    header = st.columns([3, 1, 1, 1])
    header[0].markdown("**질문**")
    header[1].markdown("**GPT OFF**")
    header[2].markdown("**GPT ON**")
    header[3].markdown("**Gemini**")
    st.markdown("<hr style='border:none;border-top:1px solid #ddd;margin:4px 0'>", unsafe_allow_html=True)

    for i, q_data in enumerate(st.session_state.questions):
        n = i + 1
        q_text = q_data.get('question', '')
        row = st.columns([3, 1, 1, 1])
        row[0].markdown(f"**Q{n}.** {q_text[:38]}{'...' if len(q_text) > 38 else ''}")
        for j, key in enumerate(ai_keys):
            ans = st.session_state.answers[key][n]
            mention = check_mention(ans, st.session_state.brand_name)
            with row[j + 1]:
                if ans.strip():
                    if mention:
                        st.markdown('<span class="mention-yes">✓ 언급</span>', unsafe_allow_html=True)
                    else:
                        st.markdown('<span class="mention-no">✗ 미언급</span>', unsafe_allow_html=True)
                else:
                    st.caption("—")

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    # Claude 분석 결과
    st.markdown("#### 🤖 Claude 분석 결과")
    st.markdown(f"""
    <div style="background:#f8f9fa;border:1px solid #e0e0e0;border-radius:10px;
                padding:20px;font-size:0.88rem;line-height:1.8;white-space:pre-wrap;">
    {st.session_state.analysis_result}
    </div>""", unsafe_allow_html=True)

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    # 크림웍스 전략 인사이트 편집
    st.markdown("#### 💜 크림웍스 전략 인사이트 편집")
    st.markdown('<div class="cw-box">보라색 영역은 크림웍스의 전략 의견입니다. 대표님이 직접 수정·보완해주세요.</div>', unsafe_allow_html=True)

    st.session_state.overall_diagnosis = st.text_area(
        "전체 현황 진단 요약",
        value=st.session_state.overall_diagnosis or st.session_state.analysis_result[:300],
        height=100,
        key="overall_diag"
    )

    for i, q_data in enumerate(st.session_state.questions):
        n = i + 1
        q_text = q_data.get('question', '')
        with st.expander(f"Q{n} 전략 인사이트 — {q_text[:40]}...", expanded=False):
            for key, label in zip(ai_keys, ai_labels_short):
                ans = st.session_state.answers[key][n]
                if ans.strip():
                    st.caption(f"**{label}:** {ans[:120]}{'...' if len(ans) > 120 else ''}")
            insight = st.text_area(
                "💜 크림웍스 전략 제안",
                value=st.session_state.cw_insights[i],
                key=f"cw_insight_{n}",
                height=90,
                placeholder="이 질문에 대한 GEO 전략 방향을 입력하세요..."
            )
            st.session_state.cw_insights[i] = insight

    st.session_state.priority_actions = st.text_area(
        "🚨 우선 실행 액션 플랜",
        value=st.session_state.priority_actions,
        height=140,
        key="priority_act"
    )

    col_back, col_next = st.columns([1, 3])
    with col_back:
        if st.button("← 답변 수정", use_container_width=True):
            st.session_state.step = 3
            st.rerun()
    with col_next:
        if st.button("📄 보고서 생성 →", type="primary", use_container_width=True):
            st.session_state.step = 5
            st.rerun()

# ─────────────────────────────────────────────
# ─────────────────────────────────────────────
# STEP 5: Word 보고서 생성
# ─────────────────────────────────────────────
elif st.session_state.step == 5:
    st.markdown('<div class="step-label">STEP 5</div>', unsafe_allow_html=True)
    st.markdown('<div class="step-title">보고서 생성</div>', unsafe_allow_html=True)

    st.markdown(f"""
    <div class="cw-box">
      💜 브랜드 컬러 <b>{st.session_state.brand_color}</b> 자동 적용 &nbsp;|&nbsp;
      크림웍스 퍼플 <b>#7030A0</b> &nbsp;|&nbsp;
      폰트 <b>페이퍼로지</b> 계열 적용
    </div>""", unsafe_allow_html=True)

    st.markdown("#### 📄 질문지")
    col_qw, col_qp = st.columns(2)

    with col_qw:
        st.markdown("**질문지 Word** — 내부 검토용")
        if st.button("📄 질문지 Word 생성", use_container_width=True):
            with st.spinner("질문지 Word 생성 중..."):
                buf = create_question_word()
            fname = f"{st.session_state.brand_name}_GEO_AI진단질문지_{datetime.now().strftime('%Y%m%d')}.docx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True, key="dl_qw"
            )

    with col_qp:
        st.markdown("**질문지 PPT** — 광고주 제안용")
        if st.button("📊 질문지 PPT 생성", use_container_width=True):
            with st.spinner("질문지 PPT 생성 중..."):
                buf = create_question_ppt()
            fname = f"{st.session_state.brand_name}_GEO_AI진단질문지_{datetime.now().strftime('%Y%m%d')}.pptx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True, key="dl_qp"
            )

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)
    st.markdown("#### 📑 전략제안서")
    col_pw, col_pp = st.columns(2)

    with col_pw:
        st.markdown("**전략제안서 Word** — 내부 검토용")
        if st.button("📄 전략제안서 Word 생성", use_container_width=True):
            with st.spinner("전략제안서 Word 생성 중..."):
                buf = create_proposal_word()
            fname = f"{st.session_state.brand_name}_GEO_전략제안서_{datetime.now().strftime('%Y%m%d')}.docx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True, key="dl_pw"
            )

    with col_pp:
        st.markdown("**전략제안서 PPT** — 광고주 제안용")
        if st.button("📊 전략제안서 PPT 생성", use_container_width=True):
            with st.spinner("전략제안서 PPT 생성 중..."):
                buf = create_proposal_ppt()
            fname = f"{st.session_state.brand_name}_GEO_전략제안서_{datetime.now().strftime('%Y%m%d')}.pptx"
            st.download_button(
                label=f"⬇️ {fname}",
                data=buf, file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True, key="dl_pp"
            )

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)

    # 전체 데이터 Excel 저장
    st.markdown("#### 💾 전체 데이터 Excel 저장")
    st.caption("브랜드 정보 + 질문 + 답변 + B2A 분석 + Claude 분석 결과를 한 파일로 저장합니다.")
    if st.button("📊 전체 데이터 Excel 저장", use_container_width=True):
        with st.spinner("Excel 파일 생성 중..."):
            buf = create_answer_excel()
        fname = f"{st.session_state.brand_name}_GEO_전체데이터_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
        st.download_button(
            label=f"⬇️ {fname} 다운로드",
            data=buf,
            file_name=fname,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
            key="dl_excel_full"
        )
        st.success("전체 데이터가 저장되었습니다!")

    st.markdown("<hr class='divider'>", unsafe_allow_html=True)
    col_back, col_new = st.columns([1, 3])
    with col_back:
        if st.button("← 인사이트 수정", use_container_width=True):
            st.session_state.step = 4
            st.rerun()
    with col_new:
        if st.button("🔄 새 브랜드 진단 시작", use_container_width=True):
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.rerun()
